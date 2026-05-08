# -*- coding: utf-8 -*-
"""
配布用ハンズオン資料 (Google Slides 取り込み用 .pptx)
Claude Chat 5日間ハンズオン研修
"""

import sys
sys.path.insert(0, ".")  # to import handson_data

from handson_data import PROFESSIONS, MULTI_PROMPTS

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor
from lxml import etree

# --- color palette ----------------------------------------------------------
PRIMARY        = RGBColor(0x6C, 0x47, 0xFF)
PRIMARY_DARK   = RGBColor(0x4A, 0x2B, 0xD4)
ACCENT         = RGBColor(0xFF, 0x6B, 0x35)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
TEXT           = RGBColor(0x1A, 0x1A, 0x2E)
MUTED          = RGBColor(0x6B, 0x72, 0x80)
SUCCESS        = RGBColor(0x16, 0xA3, 0x4A)
WARN           = RGBColor(0xFB, 0x92, 0x3C)
CARD_BG        = RGBColor(0xF8, 0xF7, 0xFC)
WARN_BG        = RGBColor(0xFF, 0xF4, 0xEE)
OK_BG          = RGBColor(0xF0, 0xFD, 0xF4)
CODE_BG        = RGBColor(0x1E, 0x1E, 0x2E)
CODE_TEXT      = RGBColor(0xE0, 0xE0, 0xF0)
CODE_LINE      = RGBColor(0x44, 0x44, 0x66)
BORDER         = RGBColor(0xE5, 0xE5, 0xEA)
SUBTLE_GREY    = RGBColor(0xFB, 0xFA, 0xFE)
DARK_BG        = RGBColor(0x1A, 0x14, 0x3D)

DAY_COLORS = [
    RGBColor(0x6C, 0x47, 0xFF),  # Day 1
    RGBColor(0x84, 0x55, 0xFF),  # Day 2
    RGBColor(0xB6, 0x52, 0xC9),  # Day 3
    RGBColor(0xE5, 0x5F, 0x8F),  # Day 4
    RGBColor(0xFF, 0x6B, 0x35),  # Day 5
]
DAY_TEXT_COLOR = WHITE  # text on day color background

HEAD_FONT = "Yu Gothic"
BODY_FONT = "Yu Gothic"
CODE_FONT = "Consolas"

# --- presentation setup -----------------------------------------------------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW_IN, SH_IN = 13.333, 7.5
BLANK = prs.slide_layouts[6]


# --- helpers ----------------------------------------------------------------
def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _kill_shadow(shape):
    sp = shape._element.spPr
    for el in sp.findall(qn("a:effectLst")):
        sp.remove(el)
    etree.SubElement(sp, qn("a:effectLst"))


def add_rect(slide, x, y, w, h, fill=None, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.line.fill.background()
    if line is not None:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    _kill_shadow(sh)
    return sh


def add_round(slide, x, y, w, h, fill=None, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.line.fill.background()
    if line is not None:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    _kill_shadow(sh)
    return sh


def add_oval(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.line.fill.background()
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    _kill_shadow(sh)
    return sh


def add_text(slide, x, y, w, h, text, *,
             size=14, bold=False, italic=False, color=TEXT,
             align="left", valign="top", font=BODY_FONT,
             char_spacing=None, line_spacing=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.word_wrap = True
    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.TOP

    lines = text.split("\n") if isinstance(text, str) else text
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = ln
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        if char_spacing:
            rPr = run._r.get_or_add_rPr()
            rPr.set("spc", str(char_spacing))
    return tb


def add_bullets_with_marker(slide, x, y, w, h, items, *, marker="▸",
                            marker_color=None, size=12, color=TEXT,
                            font=BODY_FONT, line_spacing=1.45, bold_inline=True):
    """items: list[str]. Renders left-aligned with custom marker.
    Supports inline **bold** if bold_inline=True."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.word_wrap = True
    mc = marker_color if marker_color is not None else PRIMARY
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        if marker:
            mk = p.add_run()
            mk.text = f"{marker}  "
            mk.font.name = font
            mk.font.size = Pt(size)
            mk.font.bold = True
            mk.font.color.rgb = mc
        # inline bold
        if bold_inline:
            parts = item.split("**")
            for j, part in enumerate(parts):
                if part == "":
                    continue
                run = p.add_run()
                run.text = part
                run.font.name = font
                run.font.size = Pt(size)
                run.font.bold = (j % 2 == 1)
                run.font.color.rgb = color
        else:
            run = p.add_run()
            run.text = item
            run.font.name = font
            run.font.size = Pt(size)
            run.font.color.rgb = color
    return tb


def add_step_row(slide, x, y, w, h, n, text, *, color=PRIMARY,
                 size=12, font=BODY_FONT):
    """Numbered step: circle with number + text. Returns the row shape."""
    add_rect(slide, x, y, w, h, fill=WHITE, line=BORDER)
    add_oval(slide, x + 0.18, y + (h - 0.42) / 2, 0.42, 0.42, color)
    add_text(slide, x + 0.18, y + (h - 0.42) / 2, 0.42, 0.42, str(n),
             size=11, bold=True, color=WHITE, align="center", valign="middle",
             font=HEAD_FONT)
    add_text(slide, x + 0.75, y, w - 0.85, h, text,
             size=size, color=TEXT, valign="middle", font=font, line_spacing=1.4)


def add_warn_card(slide, x, y, w, h, title, body, *, color=ACCENT, bg=WARN_BG,
                  body_size=11):
    add_round(slide, x, y, w, h, fill=bg)
    add_rect(slide, x, y, 0.08, h, fill=color)
    if title:
        add_text(slide, x + 0.30, y + 0.12, w - 0.4, 0.32, title,
                 size=12, bold=True, color=PRIMARY_DARK, font=HEAD_FONT)
        add_text(slide, x + 0.30, y + 0.45, w - 0.4, h - 0.55, body,
                 size=body_size, color=TEXT, font=BODY_FONT, line_spacing=1.5)
    else:
        add_text(slide, x + 0.30, y + 0.15, w - 0.4, h - 0.25, body,
                 size=body_size, color=TEXT, font=BODY_FONT, valign="middle",
                 line_spacing=1.5)


def add_ok_card(slide, x, y, w, h, title, body, *, body_size=11):
    add_warn_card(slide, x, y, w, h, title, body, color=SUCCESS, bg=OK_BG,
                  body_size=body_size)


def add_prompt_box(slide, x, y, w, h, text, *, size=10, font=CODE_FONT):
    """Code-style prompt box: dark background + light text."""
    add_round(slide, x, y, w, h, fill=CODE_BG)
    # left "code line" accent
    add_rect(slide, x, y + 0.15, 0.08, h - 0.30, fill=PRIMARY)
    add_text(slide, x + 0.25, y + 0.18, w - 0.4, h - 0.30, text,
             size=size, color=CODE_TEXT, font=font, line_spacing=1.4)


def add_section_header(slide, day_num, day_color):
    """Top-right corner label like 'DAY 3 / 03'."""
    add_text(slide, SW_IN - 3.0, 0.4, 2.6, 0.35,
             f"DAY {day_num}",
             size=10, bold=True, color=day_color, align="right",
             font=HEAD_FONT, char_spacing=400)


def add_footer(slide, day_num, page_no, total_pages, day_title):
    add_text(slide, 0.6, SH_IN - 0.4, 9, 0.3,
             f"DAY {day_num}：{day_title}",
             size=9, color=MUTED, font=BODY_FONT)
    add_text(slide, SW_IN - 1.5, SH_IN - 0.4, 1.0, 0.3,
             f"{page_no} / {total_pages}",
             size=9, color=MUTED, font=BODY_FONT, align="right")


# =============================================================================
# slide builders
# =============================================================================

# Day metadata
DAYS = [
    {"num": 1, "title": "Claude を自分仕様にする",
     "lead": "Projects・Memory で「毎回説明しなくていい相棒」を作る",
     "feature": "Projects / Memory"},
    {"num": 2, "title": "「読む」「まとめる」を委任する",
     "lead": "Files・Artifacts で資料処理を 1/10 の時間に",
     "feature": "Files / Artifacts"},
    {"num": 3, "title": "データとスライドを自動化する",
     "lead": "Skill (/xlsx, /pptx) で Excel職人・PPT職人を卒業",
     "feature": "Skill (/xlsx, /pptx)"},
    {"num": 4, "title": "周辺ツールと繋いで「アプリ往復」を消す",
     "lead": "Connector で Gmail・Calendar・Drive を1つの相棒に",
     "feature": "Connector (Gmail / Calendar / Drive)"},
    {"num": 5, "title": "自分の業務に組み込む",
     "lead": "学んだ機能を「明日からの定常業務」に変える",
     "feature": "棚卸し / 30日設計"},
]

# Total pages will be computed after building, so we use a placeholder for now
# and fix at the end. For simplicity, we'll write final page numbers manually.

# track current page number by day
_page_counters = {1: 0, 2: 0, 3: 0, 4: 0, 5: 0}
_day_totals = {1: 0, 2: 0, 3: 0, 4: 0, 5: 0}
_pending_footers = []  # list of (slide, day_num, page_no, day_title)


def page(day):
    _page_counters[day] += 1
    _day_totals[day] = _page_counters[day]
    return _page_counters[day]


def queue_footer(slide, day_num):
    """Queue footer to add later when totals known."""
    p = _page_counters[day_num]
    _pending_footers.append((slide, day_num, p))


def finalize_footers():
    for slide, day_num, p in _pending_footers:
        d = DAYS[day_num - 1]
        add_footer(slide, day_num, p, _day_totals[day_num], d["title"])


# -----------------------------------------------------------------------------
# Cover slide
# -----------------------------------------------------------------------------

def slide_cover():
    s = add_slide()
    set_bg(s, DARK_BG)
    # decorative ovals
    add_oval(s, -2.0, -2.5, 7.0, 7.0, RGBColor(0x6C, 0x47, 0xFF))
    add_oval(s, SW_IN - 4.5, SH_IN - 5.5, 8.0, 8.0, RGBColor(0xFF, 0x6B, 0x35))
    # darken overlay
    overlay = add_rect(s, 0, 0, SW_IN, SH_IN, fill=DARK_BG)
    sp = overlay.fill.fore_color._xFill
    a_srgb = sp.find(qn("a:srgbClr"))
    if a_srgb is not None:
        alpha = etree.SubElement(a_srgb, qn("a:alpha"))
        alpha.set("val", "78000")

    add_text(s, 1.0, 1.0, 11.5, 0.5,
             "WORKSHOP HANDOUT / 配布用ハンズオン資料",
             size=12, bold=True, color=RGBColor(0xFF, 0xB7, 0x96), font=HEAD_FONT,
             char_spacing=400)
    add_text(s, 1.0, 1.85, 11.5, 1.4,
             "Claude Chat",
             size=72, bold=True, color=WHITE, font=HEAD_FONT)
    add_text(s, 1.0, 2.95, 11.5, 1.0,
             "5日間ハンズオン研修",
             size=44, bold=True, color=RGBColor(0xFF, 0xB7, 0x96), font=HEAD_FONT)

    add_rect(s, 1.0, 4.2, 0.8, 0.06, fill=ACCENT)
    add_text(s, 1.0, 4.35, 11.5, 0.5,
             "受講者用 ／ 全Dayのハンズオン手順とプロンプトを1冊に",
             size=18, color=WHITE, font=BODY_FONT)
    add_text(s, 1.0, 4.95, 11.5, 0.5,
             "全5日 × 各2時間 ／ 合計10時間 ／ 6職種別プロンプト対応",
             size=14, color=RGBColor(0xCA, 0xC1, 0xE8), font=BODY_FONT)

    add_text(s, 1.0, SH_IN - 1.5, 11.5, 0.4,
             "教材URL： https://atsunaricoda-maker.github.io/claude-chat-course/",
             size=12, color=RGBColor(0xCA, 0xC1, 0xE8), font=BODY_FONT)
    add_text(s, 1.0, SH_IN - 1.0, 11.5, 0.4,
             "ご利用方法： Claude.ai を別タブで開き、本資料のプロンプトをコピーして送信",
             size=12, color=RGBColor(0xCA, 0xC1, 0xE8), font=BODY_FONT)
    add_text(s, 1.0, SH_IN - 0.55, 11.5, 0.3,
             "※ プロンプト内の（  ）は自分の状況に書き換えてから送信してください",
             size=10, color=RGBColor(0xA3, 0x97, 0xCC), font=BODY_FONT, italic=True)


# -----------------------------------------------------------------------------
# Section divider for each Day
# -----------------------------------------------------------------------------

def slide_day_divider(day):
    d = DAYS[day - 1]
    color = DAY_COLORS[day - 1]
    s = add_slide()
    set_bg(s, color)

    # large day number on the right
    add_text(s, SW_IN - 5.5, 0.5, 5.0, 5.0, str(day),
             size=400, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align="center",
             font=HEAD_FONT)
    # opacity overlay using XML to make number subtle
    # alternative: just trust opacity color choice. Re-draw smaller readable number:

    add_text(s, 0.7, 1.5, 11, 0.5,
             f"DAY {day}",
             size=20, bold=True, color=WHITE, font=HEAD_FONT, char_spacing=500)
    add_text(s, 0.7, 2.15, 11, 1.5, d["title"],
             size=44, bold=True, color=WHITE, font=HEAD_FONT, line_spacing=1.2)
    # divider
    add_rect(s, 0.7, 3.85, 1.0, 0.08, fill=WHITE)
    add_text(s, 0.7, 4.05, 11, 0.6, d["lead"],
             size=20, color=RGBColor(0xFF, 0xFF, 0xFF), font=BODY_FONT,
             line_spacing=1.4)
    add_text(s, 0.7, 4.85, 11, 0.4, "使う機能",
             size=11, bold=True, color=WHITE, font=HEAD_FONT, char_spacing=300)
    add_text(s, 0.7, 5.2, 11, 0.5, d["feature"],
             size=22, bold=True, color=WHITE, font=HEAD_FONT)


# -----------------------------------------------------------------------------
# 5-Day at-a-glance overview
# -----------------------------------------------------------------------------

def slide_at_a_glance():
    s = add_slide()
    set_bg(s, WHITE)
    add_text(s, SW_IN - 4.0, 0.4, 3.6, 0.35,
             "AT A GLANCE / 5日間の流れ",
             size=10, bold=True, color=PRIMARY, align="right",
             font=HEAD_FONT, char_spacing=300)
    add_text(s, 0.6, 0.45, 12.0, 0.6,
             "5日間の流れ",
             size=28, bold=True, color=TEXT, font=HEAD_FONT)
    add_text(s, 0.6, 1.05, 12.0, 0.4,
             "各Dayは独立して読めますが、Day1のProject / Memory設定が翌日以降の前提になります。",
             size=12, color=MUTED, font=BODY_FONT)

    cw = 2.46
    gap = 0.04
    bx = 0.6
    by = 1.7
    ch = 5.2

    for i, d in enumerate(DAYS):
        x = bx + i * (cw + gap)
        add_rect(s, x, by, cw, ch, fill=WHITE, line=BORDER)
        # header
        add_rect(s, x, by, cw, 0.85, fill=DAY_COLORS[i])
        add_text(s, x, by + 0.12, cw, 0.32, f"DAY {d['num']}",
                 size=14, bold=True, color=WHITE, align="center",
                 font=HEAD_FONT, char_spacing=300)
        add_text(s, x + 0.15, by + 0.45, cw - 0.3, 0.40, d["title"],
                 size=12, bold=True, color=WHITE, align="center", font=HEAD_FONT,
                 line_spacing=1.2)

        cy = by + 1.0
        # lead
        add_text(s, x + 0.20, cy, cw - 0.4, 1.2, d["lead"],
                 size=11, color=TEXT, font=BODY_FONT, line_spacing=1.5)

        # divider
        add_rect(s, x + 0.20, cy + 1.4, cw - 0.4, 0.02, fill=BORDER)

        # feature
        add_text(s, x + 0.20, cy + 1.55, cw - 0.4, 0.25, "FEATURE",
                 size=8, bold=True, color=MUTED, font=HEAD_FONT, char_spacing=200)
        add_text(s, x + 0.20, cy + 1.80, cw - 0.4, 0.4, d["feature"],
                 size=11, bold=True, color=DAY_COLORS[i], font=BODY_FONT,
                 line_spacing=1.3)

        # session count
        sessions = ["5セッション", "4セッション", "5セッション",
                    "6セッション", "5セッション"][i]
        add_text(s, x + 0.20, cy + 2.40, cw - 0.4, 0.3, "セッション数",
                 size=9, bold=True, color=MUTED, font=HEAD_FONT, char_spacing=200)
        add_text(s, x + 0.20, cy + 2.65, cw - 0.4, 0.3, sessions,
                 size=11, color=TEXT, font=BODY_FONT)

        # outcomes
        outcomes = [
            ["自分専用Project運用", "Memory基本情報設定"],
            ["PDF→自分視点要約", "1枚資料作成"],
            ["/xlsx 月次レポート", "/pptx プレゼン15分"],
            ["Gmail要約・返信", "3ツール統合"],
            ["1業務フルClaude化", "30日ロードマップ"],
        ][i]
        add_text(s, x + 0.20, cy + 3.15, cw - 0.4, 0.3, "受講後にできること",
                 size=9, bold=True, color=MUTED, font=HEAD_FONT, char_spacing=200)
        for k, oc in enumerate(outcomes):
            add_text(s, x + 0.20, cy + 3.45 + k * 0.32, cw - 0.4, 0.32,
                     f"✓ {oc}",
                     size=10, color=TEXT, font=BODY_FONT)

    # bottom
    add_round(s, 0.6, 7.0, 12.5, 0.4, fill=PRIMARY_DARK)
    add_text(s, 0.6, 7.0, 12.5, 0.4,
             "全Dayの詳細手順と全プロンプトは続きのスライドへ →",
             size=11, bold=True, color=WHITE, valign="middle", align="center",
             font=BODY_FONT)


# -----------------------------------------------------------------------------
# Generic content slide header
# -----------------------------------------------------------------------------

def header_block(slide, day_num, title, subtitle):
    color = DAY_COLORS[day_num - 1]
    # Top-right page tag
    add_text(slide, SW_IN - 3.5, 0.35, 3.0, 0.35,
             f"DAY {day_num}",
             size=10, bold=True, color=color, align="right",
             font=HEAD_FONT, char_spacing=400)
    # Title
    add_text(slide, 0.6, 0.4, 12.0, 0.7, title,
             size=26, bold=True, color=TEXT, font=HEAD_FONT, line_spacing=1.2)
    if subtitle:
        add_text(slide, 0.6, 1.05, 12.0, 0.4, subtitle,
                 size=12, color=MUTED, font=BODY_FONT)
    # Thin colored bar under header
    add_rect(slide, 0.6, 1.55, 0.8, 0.06, fill=color)


# -----------------------------------------------------------------------------
# Slide types
# -----------------------------------------------------------------------------

def make_steps_slide(day_num, title, subtitle, steps, *, warn=None, ok=None,
                     ok_at_top=False):
    """Steps + optional warn/ok card below.
    `steps`: list[str]
    `warn`: dict {title, body} or None
    `ok`: dict {title, body} or None
    """
    s = add_slide()
    set_bg(s, WHITE)
    p = page(day_num)
    color = DAY_COLORS[day_num - 1]
    header_block(s, day_num, title, subtitle)

    # body area: y = 1.85 to 6.95
    by = 1.85
    available_h = 5.1

    has_card = warn is not None or ok is not None
    if has_card:
        steps_h = available_h - 1.55
        card_y = by + steps_h + 0.15
        card_h = 1.40
    else:
        steps_h = available_h
        card_y = None
        card_h = 0

    # render steps
    n_steps = len(steps)
    if n_steps == 0:
        n_steps = 1
    row_h = min(steps_h / n_steps, 0.7)
    if row_h < 0.42:
        row_h = 0.42
    for i, st in enumerate(steps):
        y = by + i * row_h
        add_step_row(s, 0.6, y, 12.5, row_h - 0.05, i + 1, st, color=color)

    # card
    if warn:
        add_warn_card(s, 0.6, card_y, 12.5, card_h, warn.get("title"), warn["body"])
    elif ok:
        add_ok_card(s, 0.6, card_y, 12.5, card_h, ok.get("title"), ok["body"])

    queue_footer(s, day_num)
    return s


def make_cards_grid_slide(day_num, title, subtitle, cards, *, cols=2,
                          warn=None, ok=None, intro=None):
    """Grid of cards. Each card is dict {title, body (str or list[str])}.
    cards can have 2-6 items."""
    s = add_slide()
    set_bg(s, WHITE)
    p = page(day_num)
    color = DAY_COLORS[day_num - 1]
    header_block(s, day_num, title, subtitle)

    by = 1.85

    if intro:
        add_text(s, 0.6, by, 12.5, 0.4, intro,
                 size=13, color=TEXT, font=BODY_FONT, line_spacing=1.5)
        by += 0.55

    has_extra = warn is not None or ok is not None
    grid_h = (6.95 if not has_extra else 5.6) - by
    n = len(cards)
    rows = (n + cols - 1) // cols
    cw = (12.5 - (cols - 1) * 0.20) / cols
    ch = (grid_h - (rows - 1) * 0.18) / rows
    for i, card in enumerate(cards):
        r = i // cols
        c = i % cols
        x = 0.6 + c * (cw + 0.20)
        y = by + r * (ch + 0.18)
        add_round(s, x, y, cw, ch, fill=CARD_BG)
        add_rect(s, x, y, 0.08, ch, fill=color)
        add_text(s, x + 0.30, y + 0.18, cw - 0.4, 0.35, card["title"],
                 size=13, bold=True, color=PRIMARY_DARK, font=HEAD_FONT)
        if isinstance(card["body"], list):
            add_bullets_with_marker(
                s, x + 0.30, y + 0.55, cw - 0.4, ch - 0.6,
                card["body"], marker="・",
                marker_color=color, size=11, line_spacing=1.45)
        else:
            add_text(s, x + 0.30, y + 0.55, cw - 0.4, ch - 0.6,
                     card["body"],
                     size=11, color=TEXT, font=BODY_FONT, line_spacing=1.55)

    if warn:
        add_warn_card(s, 0.6, 5.7, 12.5, 1.25, warn.get("title"), warn["body"])
    elif ok:
        add_ok_card(s, 0.6, 5.7, 12.5, 1.25, ok.get("title"), ok["body"])

    queue_footer(s, day_num)
    return s


def make_flow_cards_slide(day_num, title, subtitle, flow, after_text=None,
                          cards=None, ordered_list=None):
    """Flow boxes (3 boxes connected by arrows) at top, then optional content below."""
    s = add_slide()
    set_bg(s, WHITE)
    p = page(day_num)
    color = DAY_COLORS[day_num - 1]
    header_block(s, day_num, title, subtitle)

    # flow at top
    by = 1.95
    fh = 1.5
    n = len(flow)
    arrow_w = 0.5
    box_w = (12.5 - (n - 1) * arrow_w) / n
    bx = 0.6
    for i, (label, sub) in enumerate(flow):
        x = bx + i * (box_w + arrow_w)
        if i == n - 1:
            line_color = color
            text_color = color
        else:
            line_color = BORDER
            text_color = TEXT
        add_round(s, x, by, box_w, fh, fill=WHITE, line=line_color)
        add_text(s, x + 0.2, by + 0.2, box_w - 0.4, 0.5, label,
                 size=14, bold=True, color=text_color, align="center", valign="middle",
                 font=HEAD_FONT, line_spacing=1.2)
        add_text(s, x + 0.2, by + 0.75, box_w - 0.4, 0.65, sub,
                 size=10, color=MUTED, align="center", font=BODY_FONT, line_spacing=1.4)
        if i < n - 1:
            ax = x + box_w + 0.05
            add_text(s, ax, by + (fh - 0.5) / 2, arrow_w - 0.1, 0.5, "→",
                     size=20, bold=True, color=color, align="center", valign="middle",
                     font=HEAD_FONT)

    # below content
    cy = by + fh + 0.3
    if after_text:
        add_text(s, 0.6, cy, 12.5, 0.4, after_text,
                 size=13, color=TEXT, font=BODY_FONT, line_spacing=1.5)
        cy += 0.5

    if ordered_list:
        ttl, items = ordered_list
        add_text(s, 0.6, cy, 12.5, 0.4, ttl,
                 size=14, bold=True, color=PRIMARY_DARK, font=HEAD_FONT)
        cy += 0.5
        for idx, item in enumerate(items, 1):
            add_text(s, 0.85, cy, 12.0, 0.4,
                     f"{idx}.  {item}",
                     size=12, color=TEXT, font=BODY_FONT, line_spacing=1.55)
            cy += 0.35

    if cards:
        n_cards = len(cards)
        cols = min(n_cards, 3)
        cw = (12.5 - (cols - 1) * 0.20) / cols
        ch = 1.35
        for i, card in enumerate(cards):
            x = 0.6 + (i % cols) * (cw + 0.20)
            y = cy + (i // cols) * (ch + 0.18)
            add_round(s, x, y, cw, ch, fill=CARD_BG)
            add_rect(s, x, y, 0.08, ch, fill=color)
            add_text(s, x + 0.30, y + 0.15, cw - 0.4, 0.35, card["title"],
                     size=12, bold=True, color=PRIMARY_DARK, font=HEAD_FONT)
            add_text(s, x + 0.30, y + 0.50, cw - 0.4, ch - 0.55,
                     card["body"],
                     size=10, color=TEXT, font=BODY_FONT, line_spacing=1.5)

    queue_footer(s, day_num)
    return s


def make_single_prompts_slide(day_num, title, subtitle, prompts, *,
                              after_text=None, ok=None, warn=None,
                              steps_first=None):
    """One slide with one or more 'prompt-box' style prompts.
    prompts: list of (label, text) — label can be empty.
    """
    s = add_slide()
    set_bg(s, WHITE)
    p = page(day_num)
    color = DAY_COLORS[day_num - 1]
    header_block(s, day_num, title, subtitle)

    by = 1.85
    available_h = 5.1

    if steps_first:
        steps_h = 0.42 * len(steps_first) + 0.05
        for i, st in enumerate(steps_first):
            y = by + i * 0.42
            add_step_row(s, 0.6, y, 12.5, 0.40, i + 1, st, color=color, size=11)
        by += steps_h + 0.10
        available_h -= steps_h + 0.10

    if after_text:
        add_text(s, 0.6, by, 12.5, 0.4, after_text,
                 size=12, color=TEXT, font=BODY_FONT, italic=True)
        by += 0.45
        available_h -= 0.45

    extra_card = warn is not None or ok is not None
    if extra_card:
        prompts_h = available_h - 1.20
        card_y = by + prompts_h + 0.10
        card_h = 1.05
    else:
        prompts_h = available_h
        card_y = None
        card_h = 0

    n = len(prompts)
    # share prompts_h among prompts
    each = (prompts_h - (n - 1) * 0.15) / n if n > 0 else 0
    if each < 0.6:
        each = 0.6  # minimum

    # determine font size based on text length
    for i, (label, text) in enumerate(prompts):
        py = by + i * (each + 0.15)
        # auto-size: longer texts → smaller font
        text_len = len(text)
        if text_len < 80:
            size = 12
        elif text_len < 200:
            size = 11
        elif text_len < 400:
            size = 10
        else:
            size = 9

        if label:
            add_text(s, 0.6, py - 0.05, 12.5, 0.30, label,
                     size=11, bold=True, color=PRIMARY_DARK, font=HEAD_FONT)
            py += 0.25
            each_actual = each - 0.25
        else:
            each_actual = each
        add_prompt_box(s, 0.6, py, 12.5, each_actual, text, size=size)

    if warn:
        add_warn_card(s, 0.6, card_y, 12.5, card_h, warn.get("title"), warn["body"],
                      body_size=10)
    elif ok:
        add_ok_card(s, 0.6, card_y, 12.5, card_h, ok.get("title"), ok["body"],
                    body_size=10)

    queue_footer(s, day_num)
    return s


def make_multi_prompt_short_slide(day_num, title, subtitle, prompt_key):
    """Single slide showing all 6 profession variants in 2x3 grid (for short/medium)."""
    s = add_slide()
    set_bg(s, WHITE)
    p = page(day_num)
    color = DAY_COLORS[day_num - 1]
    header_block(s, day_num, title, subtitle)

    add_text(s, 0.6, 1.85, 12.5, 0.35,
             "あなたの職種のセクションをコピーして使ってください（6職種版）",
             size=11, color=MUTED, font=BODY_FONT, italic=True)

    by = 2.20
    available_h = 4.85
    cols = 3
    rows = 2
    cw = (12.5 - (cols - 1) * 0.18) / cols
    ch = (available_h - (rows - 1) * 0.18) / rows

    prompts = MULTI_PROMPTS[prompt_key]
    # Determine font size based on max length
    max_len = max(len(prompts[k]) for k, _ in PROFESSIONS)
    if max_len < 200:
        size = 9
    elif max_len < 400:
        size = 8
    else:
        size = 7

    for i, (key, label) in enumerate(PROFESSIONS):
        c = i % cols
        r = i // cols
        x = 0.6 + c * (cw + 0.18)
        y = by + r * (ch + 0.18)
        # outer card
        add_round(s, x, y, cw, ch, fill=CODE_BG)
        # profession label (top bar)
        add_rect(s, x, y, cw, 0.40, fill=color)
        add_text(s, x + 0.15, y, cw - 0.30, 0.40, label,
                 size=11, bold=True, color=WHITE, valign="middle",
                 font=HEAD_FONT, char_spacing=200)
        # prompt content
        text = prompts[key]
        add_text(s, x + 0.20, y + 0.50, cw - 0.40, ch - 0.60, text,
                 size=size, color=CODE_TEXT, font=CODE_FONT, line_spacing=1.35)

    # footer note
    add_text(s, 0.6, 7.10, 12.5, 0.30,
             f"プロンプト：{prompt_key}  ／  全文は教材URLから検索可能",
             size=9, color=MUTED, font=BODY_FONT, italic=True)

    queue_footer(s, day_num)
    return s


def make_multi_prompt_long_slide(day_num, title, subtitle, prompt_key, group):
    """For LONG multi-prompts: split into 2 slides.
    group: 'a' = first 3 (sales/marketing/hr), 'b' = last 3 (finance/admin/manager)
    Each slide shows 3 variants in 1 row × 3 cols.
    """
    s = add_slide()
    set_bg(s, WHITE)
    p = page(day_num)
    color = DAY_COLORS[day_num - 1]
    header_block(s, day_num, title, subtitle)

    if group == "a":
        idx_range = (0, 3)
        sub_note = "営業 / マーケ・企画 / 人事 版（次のページに 経理 / 総務・事務 / 管理職 版あり）"
    else:
        idx_range = (3, 6)
        sub_note = "経理 / 総務・事務 / 管理職 版（前のページに 営業 / マーケ・企画 / 人事 版あり）"

    add_text(s, 0.6, 1.85, 12.5, 0.35, sub_note,
             size=11, color=MUTED, font=BODY_FONT, italic=True)

    by = 2.20
    available_h = 4.85
    cols = 3
    cw = (12.5 - (cols - 1) * 0.18) / cols
    ch = available_h

    prompts = MULTI_PROMPTS[prompt_key]
    keys_subset = [PROFESSIONS[i] for i in range(*idx_range)]
    max_len = max(len(prompts[k]) for k, _ in keys_subset)
    if max_len < 600:
        size = 8.5
    elif max_len < 800:
        size = 7.5
    else:
        size = 7

    for i, (key, label) in enumerate(keys_subset):
        x = 0.6 + i * (cw + 0.18)
        y = by
        add_round(s, x, y, cw, ch, fill=CODE_BG)
        add_rect(s, x, y, cw, 0.40, fill=color)
        add_text(s, x + 0.15, y, cw - 0.30, 0.40, label,
                 size=11, bold=True, color=WHITE, valign="middle",
                 font=HEAD_FONT, char_spacing=200)
        text = prompts[key]
        add_text(s, x + 0.20, y + 0.50, cw - 0.40, ch - 0.60, text,
                 size=size, color=CODE_TEXT, font=CODE_FONT, line_spacing=1.35)

    add_text(s, 0.6, 7.10, 12.5, 0.30,
             f"プロンプト：{prompt_key} ({group})  ／  全文は教材URLから検索可能",
             size=9, color=MUTED, font=BODY_FONT, italic=True)

    queue_footer(s, day_num)
    return s


def make_checklist_slide(day_num, title, subtitle, items, *, ok=None, warn=None,
                         link=None):
    """Checklist + optional ok/warn card + optional link box."""
    s = add_slide()
    set_bg(s, WHITE)
    p = page(day_num)
    color = DAY_COLORS[day_num - 1]
    header_block(s, day_num, title, subtitle)

    by = 1.85

    if ok or warn:
        list_h = 3.6
    else:
        list_h = 4.5
    if link:
        list_h -= 0.7

    # checklist
    for i, item in enumerate(items):
        y = by + i * 0.42
        add_oval(s, 0.6, y + 0.06, 0.30, 0.30, WHITE, line=color)
        add_rect(s, 0.6, y + 0.06, 0.30, 0.30, fill=WHITE, line=color)
        add_text(s, 1.1, y, 11.7, 0.42,
                 item, size=12, color=TEXT, valign="middle", font=BODY_FONT,
                 line_spacing=1.4)
    by += len(items) * 0.42 + 0.20

    # ok/warn card
    if ok:
        add_ok_card(s, 0.6, by, 12.5, 1.20, ok.get("title"), ok["body"], body_size=11)
        by += 1.30
    elif warn:
        add_warn_card(s, 0.6, by, 12.5, 1.20, warn.get("title"), warn["body"], body_size=11)
        by += 1.30

    if link:
        add_round(s, 4.0, by, 5.3, 0.6, fill=color)
        add_text(s, 4.0, by, 5.3, 0.6, link,
                 size=14, bold=True, color=WHITE, align="center", valign="middle",
                 font=HEAD_FONT)

    queue_footer(s, day_num)
    return s


def make_text_only_slide(day_num, title, subtitle, content_lines):
    """Generic title + body text slide."""
    s = add_slide()
    set_bg(s, WHITE)
    p = page(day_num)
    color = DAY_COLORS[day_num - 1]
    header_block(s, day_num, title, subtitle)

    by = 1.85
    add_bullets_with_marker(s, 0.6, by, 12.5, 4.5, content_lines,
                            marker="▸", marker_color=color, size=14,
                            line_spacing=1.65)
    queue_footer(s, day_num)
    return s


# =============================================================================
# Closing slide
# =============================================================================

def slide_closing():
    s = add_slide()
    set_bg(s, DARK_BG)
    add_oval(s, SW_IN - 4.5, -2.5, 7.0, 7.0, RGBColor(0xFF, 0x6B, 0x35))
    add_oval(s, -3.0, SH_IN - 5.0, 8.0, 8.0, RGBColor(0x6C, 0x47, 0xFF))
    overlay = add_rect(s, 0, 0, SW_IN, SH_IN, fill=DARK_BG)
    sp = overlay.fill.fore_color._xFill
    a_srgb = sp.find(qn("a:srgbClr"))
    if a_srgb is not None:
        alpha = etree.SubElement(a_srgb, qn("a:alpha"))
        alpha.set("val", "78000")

    add_text(s, 1.0, 1.7, 11.5, 0.5,
             "CONGRATULATIONS / 修了",
             size=12, bold=True, color=RGBColor(0xFF, 0xB7, 0x96), font=HEAD_FONT,
             char_spacing=400)
    add_text(s, 1.0, 2.3, 11.5, 1.4,
             "5日間お疲れさまでした",
             size=42, bold=True, color=WHITE, font=HEAD_FONT, line_spacing=1.2)
    add_text(s, 1.0, 3.4, 11.5, 1.4,
             "明日からの仕事が、今日までと違うことを願っています。",
             size=22, bold=True, color=RGBColor(0xFF, 0xB7, 0x96), font=HEAD_FONT)

    add_text(s, 1.0, 4.7, 11.5, 1.0,
             "困ったときは、いちばん最初に作った Project に「困っている」と話しかけることから始めてください。",
             size=14, color=RGBColor(0xCA, 0xC1, 0xE8), font=BODY_FONT, line_spacing=1.7)

    add_round(s, 1.0, 5.85, 11.3, 1.2, fill=WHITE)
    add_text(s, 1.4, 5.95, 10.5, 0.4, "教材は本日からアクセス可能です",
             size=12, bold=True, color=PRIMARY_DARK, font=HEAD_FONT)
    add_text(s, 1.4, 6.3, 10.5, 0.4,
             "https://atsunaricoda-maker.github.io/claude-chat-course/",
             size=18, bold=True, color=PRIMARY_DARK, font=HEAD_FONT)
    add_text(s, 1.4, 6.65, 10.5, 0.35,
             "30日定着ロードマップを実行 ／ 困ったら万能プロンプト集 ／ 同僚に3分デモ",
             size=11, color=MUTED, font=BODY_FONT)


# =============================================================================
# main build flow
# =============================================================================

slide_cover()
slide_at_a_glance()


# -----------------------------------------------------------------------------
# DAY 1 — Claude を自分仕様にする
# -----------------------------------------------------------------------------
slide_day_divider(1)

# 1.1 アカウント準備
make_steps_slide(
    day_num=1,
    title="① アカウント準備",
    subtitle="所要 5分 ／ 既にログイン済みの方は次のスライドへ",
    steps=[
        "ブラウザで claude.ai を開く（Chrome か Edge 推奨）",
        "画面中央の Sign up（新規登録）または Sign in（ログイン）をクリック",
        "Google アカウント or メールアドレスでログイン",
        "画面右上に自分のアイコンが表示されればログイン完了",
    ],
    warn={"title": "会社で使う場合の注意",
          "body": "会社のメールアドレスでアカウントを作る前に、情シスに「Claude.ai を使ってよいか／どのプランか」を確認してください。Free / Pro / Team / Enterprise でデータ取扱いが異なります。"},
)

# 1.2 画面レイアウトを覚える
make_cards_grid_slide(
    day_num=1,
    title="② 画面レイアウトを覚える",
    subtitle="どこに何があるかを把握すると迷わなくなる",
    cards=[
        {"title": "左サイドバー",
         "body": ["新しいチャット（左上 +）",
                  "**Chats**：過去の会話一覧",
                  "**Projects**：今日の主役",
                  "**Artifacts**：作った成果物"]},
        {"title": "中央 / 右",
         "body": ["真ん中：会話エリア",
                  "下部：入力ボックス（📎で添付）",
                  "入力欄左下：**モデル切替**",
                  "入力欄右下：**送信ボタン**"]},
    ],
    ok={"title": "モデルの使い分け（迷ったら Sonnet）",
        "body": "Opus：難しい思考・長文 ／ Sonnet：日常業務の万能 ／ Haiku：とにかく速く"},
)

# 1.3 Project とは何か
make_flow_cards_slide(
    day_num=1,
    title="③ Project とは何か",
    subtitle="「文脈付きの専用チャットルーム」",
    flow=[
        ("通常チャット", "毎回ゼロから"),
        ("Project", "背景＋資料を共有"),
        ("毎回\n自分仕様の返答", ""),
    ],
    ordered_list=("Projectで設定する2つのもの", [
        "**System instructions（指示）**：自分の役割・出力ルール",
        "**Project knowledge（資料）**：参照させたいファイル",
    ]),
    after_text="一度入れたら、その Project 内のすべてのチャットで自動的に使われる",
)

# 1.4 ハンズオン①：Project を作る
make_steps_slide(
    day_num=1,
    title="ハンズオン①：Project を作る",
    subtitle="所要 10分 ／ Lv1：自分専用の文体ルールを持たせる",
    steps=[
        "左サイドバーの Projects（プロジェクト）をクリック",
        "右上の + Create project（プロジェクトを作成）をクリック",
        "Name に「マイ業務アシスタント」と入力",
        "Description は空欄で OK → Create project",
        "右側の Set custom instructions → 次のスライドのプロンプトを貼り付け",
    ],
)

# 1.5 コピペ①：Project の指示文 (LONG multi)
make_multi_prompt_long_slide(
    day_num=1,
    title="コピペ①：Project の指示文",
    subtitle="あなたの職種のセクションをコピー → Custom instructions 欄に貼り付け",
    prompt_key="project_instructions",
    group="a",
)
make_multi_prompt_long_slide(
    day_num=1,
    title="コピペ①：Project の指示文（続き）",
    subtitle="あなたの職種のセクションをコピー → Custom instructions 欄に貼り付け",
    prompt_key="project_instructions",
    group="b",
)

# 1.6 動作確認 (SHORT multi)
make_multi_prompt_short_slide(
    day_num=1,
    title="動作確認：いま設定した指示が効くか試す",
    subtitle="Project 内で新規チャット → 自分の職種版を送信  ／  こうなれば成功：質問してくる・結論先・丁寧体",
    prompt_key="verify_choei",
)

# 1.7 ハンズオン②：自分の文体を学習
make_steps_slide(
    day_num=1,
    title="ハンズオン②：自分の文体を学習させる（Lv2）",
    subtitle="所要 25分 ／ 過去資料を Project に食べさせる",
    steps=[
        "用意：自分が過去に書いたメール 5通 ＋ 報告書/提案書 2本（個人情報は黒塗り）",
        "Project 画面右側 Project knowledge エリアの + Add content をクリック",
        "ファイルをドラッグ＆ドロップ（PDF / Word / テキスト OK）",
        "Custom instructions に下の追記文（次スライド）を加える",
        "新規チャットで実際にメール下書きを依頼してみる",
    ],
)

# 1.8 コピペ②：文体学習の追記文
make_single_prompts_slide(
    day_num=1,
    title="コピペ②：文体学習の追記文 ／ 試してみるプロンプト",
    subtitle="既存の Custom instructions の末尾に追記 → 新規チャットで2つ目を試す",
    prompts=[
        ("【追記文】Custom instructions に追加",
         "【私の文体について】\n"
         "Project knowledge にアップロードされている過去の文書は、すべて私が実際に書いたものです。\n"
         "私が文章作成を依頼したときは、これらの文書から以下を学習し踏襲してください：\n\n"
         "・改行のタイミングと段落の長さ\n"
         "・よく使う言い回し・接続詞\n"
         "・敬語の硬さの度合い\n"
         "・箇条書きと文章のバランス\n"
         "・締めくくり方の癖\n\n"
         "ただし、過去文書に明らかな誤字脱字や論理破綻があった場合は、それは真似しないでください。"),
        ("【試す】新規チャットで送ってみる",
         "来週水曜14:00からの定例会を、別件のため15:00開始に変更したいというメールを、"
         "参加者5名（社内）に送りたい。私の文体で下書きを作って。"),
    ],
)

# 1.9 Memory を有効化する
make_steps_slide(
    day_num=1,
    title="④ Memory を有効化する",
    subtitle="所要 15分 ／ 別のチャットでも記憶される仕組み",
    steps=[
        "左下の自分のアイコン → Settings（設定）",
        "左メニュー Profile（プロフィール）または Personalization（パーソナライズ）",
        "Memory（メモリ）または Reference recent chats（最近のチャットを参照）をオン",
        "設定を閉じる → 通常のチャット（Project外でも可）を開く",
    ],
    warn={"title": "Memory と Project の違い",
          "body": "**Project**：そのProject内だけで使われる。仕事の文脈ごとに分けたい時。\n**Memory**：すべてのチャットで使われる。あなた個人の前提情報。"},
)

# 1.10 Memory に覚えさせる基本情報 (MEDIUM multi)
make_multi_prompt_short_slide(
    day_num=1,
    title="コピペ③：Memory に覚えさせる基本情報",
    subtitle="普通のチャット（Project外）で、あなたの職種版を送信",
    prompt_key="memory_basic",
)

# 1.11 本日の宿題
make_checklist_slide(
    day_num=1,
    title="本日の宿題（次回までに）",
    subtitle="明日の Day 2 を効果的にするために",
    items=[
        "Project に、追加で 3〜5通の自分のメール をアップロード",
        "Memory に、思いついた「いつも説明していること」を1つ追加",
        "明日のために、要約してみたい PDF（10ページ以上）を1つ手元に用意",
        "明日のために、文字起こし or 録音した会議の議事録を1つ用意（30分以上）",
    ],
    ok={"title": "今日の達成",
        "body": "・自分仕様のProjectを作成   ・自分の文体を学習させた   ・Memoryで個人情報を記憶させた"},
    link="Day 2 へ進む  →",
)


# -----------------------------------------------------------------------------
# DAY 2 — 「読む」「まとめる」を委任する
# -----------------------------------------------------------------------------
slide_day_divider(2)

# 2.1 ファイル投入の基本
make_steps_slide(
    day_num=2,
    title="ファイル投入の基本",
    subtitle="どこから何を入れるか",
    steps=[
        "入力ボックスの左下 📎（クリップ）アイコンをクリック",
        "ファイルを選ぶ（PDF / Word / Excel / 画像 / 音声 / テキスト）",
        "複数選択も可。最大数十MBまで（プランによる）",
        "添付されたら入力欄にプロンプトを書いて送信",
    ],
    warn={"title": "機密情報の判断軸",
          "body": "個人情報（氏名・連絡先・口座番号）→ 必ず黒塗り or 削除  ／  顧客名・取引先名 → 「A社」「B社」に置換  ／  未公開の財務情報 → 会社のClaudeポリシーに従う"},
)

# 2.2 ハンズオン①：PDF を要約 (Lv1)
make_single_prompts_slide(
    day_num=2,
    title="ハンズオン①：PDFを要約（Lv1）",
    subtitle="所要 25分 ／ 用意したPDFを使う",
    steps_first=[
        "Day1で作った「マイ業務アシスタント」Project を開く",
        "+ New chat で新規チャット",
        "📎 でPDFを添付",
        "下のプロンプトを送信",
    ],
    prompts=[
        ("",
         "添付したPDFを、私（の職種）の視点で読み解いてください。\n\n"
         "以下の構成で出力してください：\n\n"
         "1. 【1行サマリー】この文書を一言で\n"
         "2. 【私が知っておくべきこと TOP5】重要度順に。各項目の根拠ページも明記\n"
         "3. 【私の業務への影響】3つの観点で（業務プロセス／関係者／意思決定）\n"
         "4. 【次に取るべきアクション】3つ。優先度A/B/Cも付けて\n"
         "5. 【素人への説明】この文書を知らない同僚に5行で説明する原稿\n\n"
         "不明点があれば最後にまとめて質問してください。"),
    ],
)

# 2.3 応用：質問を重ねて深掘りする
make_single_prompts_slide(
    day_num=2,
    title="応用：質問を重ねて深掘りする",
    subtitle="1回の返答で終わらせず、対話で価値を引き出す ／ 続けて使えるプロンプト集",
    prompts=[
        ("① さらに具体例を",
         "TOP5の3番について、もっと詳しく、具体例を3つ挙げて説明して。"),
        ("② 報告メールに転用",
         "このPDFの内容を、上司への報告メール（300字）にして。私の文体で。"),
        ("③ 違和感の指摘",
         "このPDFと一般的な業界の常識を比較して、違和感のある主張があれば指摘して。"),
    ],
    ok={"title": "体感ポイント",
        "body": "1つの資料を、要約・批評・転用の3方向で使えるのが Claude の強み。「読む時間」が「使う時間」に変わります。"},
)

# 2.4 ハンズオン②：議事録→ToDo
make_single_prompts_slide(
    day_num=2,
    title="ハンズオン②：議事録→ToDo（Lv2）",
    subtitle="所要 35分 ／ 会議録音 or 文字起こしを使う",
    steps_first=[
        "新規チャットを開く",
        "📎 で議事録ファイル添付（音声直接でも OK：.mp3 .m4a .wav）",
        "下のプロンプトを送る",
    ],
    prompts=[
        ("",
         "添付は会議の議事録（または音声）です。以下を順に行ってください。\n\n"
         "【ステップ1：構造化】\n"
         "- 会議の目的を1行で\n"
         "- 参加者を一覧化（発言から推測）\n"
         "- 議題ごとに分けて要点を箇条書き\n\n"
         "【ステップ2：意思決定の抽出】\n"
         "表で整理。列は：決定事項 / 決定者 / 理由 / 反対意見の有無\n\n"
         "【ステップ3：ToDoリスト】\n"
         "表で整理。列は：タスク / 担当者 / 期限 / 優先度 / 関連議題\n\n"
         "【ステップ4：未決事項】\n"
         "未だ決まっていない議題と、次回までに誰が何を準備すべきかを整理\n\n"
         "【ステップ5：共有メール下書き】\n"
         "参加者全員に送る議事録共有メールを、私の文体で作成。\n"
         "（件名・本文・末尾の「ご質問あればこのメールに返信ください」まで含めて）"),
    ],
)

# 2.5 議事録ハンズオンの仕上げ
make_single_prompts_slide(
    day_num=2,
    title="議事録ハンズオンの仕上げ",
    subtitle="使える成果物に整える ／ 追加で送るプロンプト",
    prompts=[
        ("① スプレッドシート向けに整形",
         "ToDoリストをコピペでスプレッドシートに貼れる形式（タブ区切り）にして。"),
        ("② 抜けた論点を提示",
         "この会議で「もし私だったら追加で議論すべきだった論点」を3つ提示して。理由も。"),
    ],
    ok={"title": "体感ポイント",
        "body": "1時間の会議の「後処理90分」が10分に。さらに「次の会議の質を上げる視点」までもらえる。"},
)

# 2.6 ハンズオン③：Artifacts で1枚資料 (Lv1)
make_steps_slide(
    day_num=2,
    title="ハンズオン③：Artifacts で1枚資料（Lv1）",
    subtitle="所要 30分 ／ 「白紙からの脱出」  ／  Artifacts ＝ 成果物が右側に独立パネルで表示",
    steps=[
        "新規チャットで次のスライドのプロンプトを送る",
        "右側にプレビューが現れることを確認",
        "「色を変えて」「ロゴを左上に」など追加指示で編集",
        "右上の Copy（コピー）または Download（ダウンロード）で取り出す",
    ],
)

# 2.7 コピペ：1枚資料プロンプト (LONG multi → 2 slides)
make_multi_prompt_long_slide(
    day_num=2,
    title="コピペ：1枚資料プロンプト",
    subtitle="あなたの職種版をコピー（テーマは職種別に最適化済み）",
    prompt_key="one_page_summary",
    group="a",
)
make_multi_prompt_long_slide(
    day_num=2,
    title="コピペ：1枚資料プロンプト（続き）",
    subtitle="あなたの職種版をコピー（テーマは職種別に最適化済み）",
    prompt_key="one_page_summary",
    group="b",
)

# 2.8 Artifacts を編集で育てる
make_single_prompts_slide(
    day_num=2,
    title="Artifacts を編集で育てる（Lv2）",
    subtitle="所要 15分 ／ 「会話で育てる」感覚を体得  ／  連続で送ってみるプロンプト",
    prompts=[
        ("① 色変更",
         "真ん中のカードの背景色を薄いオレンジに変えて、目立たせて。"),
        ("② 追記",
         "下部の結論枠に「3ヶ月後の見直しポイント」を追加して。"),
        ("③ 平易な版を作成",
         "このサマリーをそのまま社外向けに使えるよう、専門用語を平易な言葉に置き換えたバージョンも作って。"),
    ],
    ok={"title": "気づき",
        "body": "普通の文書なら作り直しになる修正が、会話1往復で済む。「初稿を作る労力」と「修正する労力」の両方が減る。"},
)

# 2.9 本日の宿題
make_checklist_slide(
    day_num=2,
    title="本日の宿題",
    subtitle="明日の Day 3（データ・スライド自動化）に備える",
    items=[
        "明日に向けて、Excel or CSVデータを1つ用意（売上・経費・アンケート集計など）",
        "明日に向けて、PowerPointで作りたいテーマを1つ決める",
        "今日作った Artifacts を実際に印刷して同僚に見せ、反応を見る",
        "議事録ハンズオンの結果を、関係者に共有メールとして実際に送ってみる",
    ],
    ok={"title": "今日の達成",
        "body": "・PDFを「自分の視点」で要約できた  ・議事録を実用的なToDoリストに変換  ・Artifactsで1枚資料を作って編集できた"},
    link="Day 3 へ進む  →",
)


# -----------------------------------------------------------------------------
# DAY 3 — データとスライドを自動化する
# -----------------------------------------------------------------------------
slide_day_divider(3)

# 3.1 Skill とは何か
make_flow_cards_slide(
    day_num=3,
    title="Skill とは何か",
    subtitle="「Claudeに専門のプロ知識を装着させるパッケージ」",
    flow=[
        ("白紙のClaude", "毎回ゼロから指示"),
        ("+ Skill", "「○○職人」を装着"),
        ("毎回プロ品質", "属人化ゼロ"),
    ],
    cards=[
        {"title": "📄 文書作成系",
         "body": "/xlsx /pptx /docx /pdf — Excel・PPT・Word生成、PDF処理"},
        {"title": "🎨 デザイン系",
         "body": "/brand-guidelines /canvas-design /theme-factory — ブランド準拠・ポスター・配色"},
        {"title": "💬 コミュニケーション系",
         "body": "/internal-comms /doc-coauthoring — 社内文書フォーマット・共著ワークフロー"},
        {"title": "⚙ 業務効率化系",
         "body": "/schedule /loop /skill-creator — 定期実行・繰り返し・独自Skill作成"},
        {"title": "✓ 品質の安定",
         "body": "誰が呼んでも同じプロ品質で出力される"},
        {"title": "✓ 時間短縮",
         "body": "長いプロンプトでなく / で1コマンド"},
    ],
)

# 3.2 Skill の設定方法
make_steps_slide(
    day_num=3,
    title="Skill の設定方法（事前準備）",
    subtitle="/ で呼び出す前に、まず有効化が必要です",
    steps=[
        "画面左下の自分のアイコン → Settings（設定）",
        "左メニューから Capabilities（機能）または Tools（ツール）を選ぶ",
        "File creation または Skills のセクションを探す",
        "使いたいものを ON に（例：Excel生成 / PowerPoint生成 / PDF処理）",
        "設定を閉じる → 入力欄で / を打ち、候補リストに出れば成功",
    ],
    warn={"title": "会社のClaudeを使う場合の注意",
          "body": "Skill の一部は管理者（情シス・ワークスペース管理者）でしか有効化できないことがあります。"
                  "設定画面で項目が見えない or グレーアウトしている場合は管理者に「○○のSkillを使いたい」と相談してください。"},
)

# 3.3 Skill の呼び出し方
make_steps_slide(
    day_num=3,
    title="Skill の呼び出し方",
    subtitle="/（半角スラッシュ）を打つだけ",
    steps=[
        "入力欄に / （半角スラッシュ）を打つ",
        "候補リストが出るので、使いたい Skill を選ぶ（例：/xlsx）",
        "続けて指示を書いて送信",
    ],
    warn={"title": None,
          "body": "もしリストに /xlsx が見えない場合は、Settings → Capabilities や Tools でファイル作成系の機能をオンにしてください。"},
)

# 3.4 ハンズオン①：/xlsx で表を作る (Lv1) — steps slide
make_steps_slide(
    day_num=3,
    title="ハンズオン①：/xlsx で表を作る（Lv1）",
    subtitle="所要 30分 ／ 用意したCSVを使う ／ 次のスライドにプロンプト",
    steps=[
        "新規チャットを開く（Project外で OK）",
        "📎 でCSV / Excel を添付",
        "次のスライドのプロンプトを送信（職種別6種）",
        "完成した .xlsx を Download（ダウンロード）で取得",
    ],
)

# 3.4b xlsx_monthly multi prompt
make_multi_prompt_short_slide(
    day_num=3,
    title="コピペ：/xlsx 月次レポートプロンプト",
    subtitle="あなたの職種版をコピー（KPI構成は職種別に最適化済み）",
    prompt_key="xlsx_monthly",
)

# 3.5 ハンズオン②：グラフと分析 (Lv2)
make_single_prompts_slide(
    day_num=3,
    title="ハンズオン②：グラフと分析（Lv2）",
    subtitle="所要 30分 ／ 同じデータの続き、または別データで",
    prompts=[
        ("① 分析レポート",
         "/xlsx 同じデータを使って、意思決定に使える分析レポートを作ってください。\n\n"
         "【含めてほしいもの】\n"
         "1. 月別推移の折れ線グラフ（複数項目を重ねる）\n"
         "2. カテゴリ別の構成比（円グラフ または 100%積み上げ棒）\n"
         "3. 異常値の自動検出（直近平均±2σを超えるものに色付け）\n"
         "4. 前年同月比の表\n"
         "5. 「数字が伝える3つのメッセージ」というセル（要約テキスト）\n\n"
         "【追加で】\n"
         "- グラフは別シート「グラフ」にまとめる\n"
         "- 「メッセージ」セルは私が経営会議で読み上げられる文体で\n\n"
         "不足するデータがあれば、何が必要かを先に教えてください。"),
        ("② 想定質問",
         "このグラフを見て、私が経営会議で受けそうな質問TOP5と、それぞれへの想定回答を準備して。"),
    ],
)

# 3.6 ハンズオン③：/pptx でスライド (Lv1) — steps
make_steps_slide(
    day_num=3,
    title="ハンズオン③：/pptx でスライド（Lv1）",
    subtitle="所要 25分 ／ 議事録 or レポートからスライド化 ／ 次のスライドにプロンプト",
    steps=[
        "新規チャット → 元ネタとなる文書（議事録、レポート、メモ）を 📎 で添付",
        "次のスライドのプロンプトを送信（職種別6種）",
        "完成した .pptx を Download（ダウンロード）",
        "PowerPointで開いて、必要に応じて手で微修正",
    ],
)

# 3.6b pptx_15slide multi prompt
make_multi_prompt_short_slide(
    day_num=3,
    title="コピペ：/pptx 15枚スライドプロンプト",
    subtitle="あなたの職種版をコピー（テーマ・構成は職種別に最適化済み）",
    prompt_key="pptx_15slide",
)

# 3.7 ハンズオン④：ブランド準拠 (Lv2)
make_single_prompts_slide(
    day_num=3,
    title="ハンズオン④：ブランド準拠（Lv2）",
    subtitle="所要 20分 ／ 会社の見た目に合わせる",
    prompts=[
        ("",
         "/pptx さきほど作ったプレゼンを、以下のブランド設定で作り直してください。\n\n"
         "【カラーパレット】\n"
         "- メインカラー：（例：#003366）\n"
         "- アクセントカラー：（例：#FF6600）\n"
         "- 背景：白\n"
         "- 文字：濃いグレー（#333333）\n\n"
         "【フォント】\n"
         "- 見出し：游ゴシック Bold\n"
         "- 本文：游ゴシック Regular\n"
         "- 数字を強調する箇所：Bebas Neue or Impact\n\n"
         "【レイアウト】\n"
         "- 各スライド左上に細い色帯（メインカラー）\n"
         "- 表紙のみ全面メインカラー、文字白\n"
         "- ページ番号を右下、小さく\n\n"
         "【ロゴ】\n"
         "- ロゴ画像は後で差し込むので、表紙とすべてのスライド右下に「[LOGO]」プレースホルダーを置く"),
    ],
    ok={"title": "応用",
        "body": "会社のスライドテンプレートを 📎 で添付して「このテンプレートに沿って作って」と頼むと、より精度が上がります。"},
)

# 3.8 本日の宿題
make_checklist_slide(
    day_num=3,
    title="本日の宿題",
    subtitle="明日の Day 4（ツール連携）に備える",
    items=[
        "作った .xlsx を実際に上司や同僚に見せて、フィードバックをもらう",
        "作った .pptx を、明日（または来週）の打ち合わせで実際に使う",
        "明日のために、Gmail（または会社で使うメールアプリ）にログインできることを確認",
        "明日のために、Googleカレンダー（または会社のカレンダー）にもログインしておく",
    ],
    warn={"title": "明日の前の確認",
          "body": "会社のセキュリティポリシーで、Claude に業務メール／カレンダーを接続してよいか必ず確認してください。"},
    link="Day 4 へ進む  →",
)


# -----------------------------------------------------------------------------
# DAY 4 — 周辺ツールと繋いで「アプリ往復」を消す
# -----------------------------------------------------------------------------
slide_day_divider(4)

# 4.1 Connector とは何か
make_flow_cards_slide(
    day_num=4,
    title="Connector とは何か",
    subtitle="「外部ツールと Claude を直接つなぐ橋」",
    flow=[
        ("単独のClaude", "あなたが手動でコピペ"),
        ("+ Connector", "Claudeが直接アクセス"),
        ("アプリ往復ゼロ", "横断検索・統合分析"),
    ],
    cards=[
        {"title": "📧 メール・カレンダー",
         "body": "Gmail / Outlook — 受信箱要約・返信下書き  ／  Google Calendar — 予定整理・会議準備"},
        {"title": "📁 ドキュメント・ストレージ",
         "body": "Google Drive — Docs/Sheets/Slides横断検索  ／  OneDrive / Dropbox / Notion"},
        {"title": "💬 チームコミュニケーション",
         "body": "Slack — チャネル要約・履歴検索  ／  Microsoft Teams"},
        {"title": "🛠 業務システム",
         "body": "Linear / Asana — タスク管理  ／  Salesforce / HubSpot — CRM  ／  GitHub / GitLab"},
        {"title": "✓ アプリ往復ゼロ",
         "body": "タブを開きまくる作業が消える"},
        {"title": "✓ 横断的な情報統合",
         "body": "Gmail × Calendar × Drive を一括で参照"},
    ],
)

# 4.2 Connector の設定方法
make_steps_slide(
    day_num=4,
    title="Connector の設定方法",
    subtitle="接続の場所・手順・注意点",
    steps=[
        "画面左下の自分のアイコン → Settings（設定）",
        "左メニュー Connectors（コネクタ）をクリック",
        "利用可能なサービス一覧から接続したいものを選択",
        "Connect（接続）をクリック → 各サービスのログイン画面に遷移",
        "権限を確認して許可（読み取りのみが最も安全）",
    ],
    warn={"title": "接続前に必ず確認",
          "body": "・会社の情シスポリシーでこのConnectorが許可されているか\n・個人アカウントと業務アカウントを混同していないか\n・付与する権限の範囲（読み取りのみ／書き込みも）"},
)

# 4.3 ハンズオン①：Gmail 接続と要約
make_single_prompts_slide(
    day_num=4,
    title="ハンズオン①：Gmail 接続と要約",
    subtitle="所要 30分",
    steps_first=[
        "Settings → Connectors → Gmail の Connect（接続）",
        "Googleの認証画面で「読み取り権限のみ」推奨でログイン",
        "新規チャットを開いて下のプロンプトを送信",
    ],
    prompts=[
        ("",
         "Gmailで、過去24時間に届いた未読メールを取得して、以下の形式でまとめてください。\n\n"
         "【出力形式】\n"
         "表で整理。列は：\n"
         "- 重要度（A=即対応 / B=今日中 / C=今週中 / D=参考のみ）\n"
         "- 送信者\n"
         "- 件名\n"
         "- 30字以内の要約\n"
         "- 私が取るべきアクション\n"
         "- 推奨返信ドラフト（必要なものだけ・3行以内）\n\n"
         "【判定基準】\n"
         "- 上司・顧客からの直接の問いかけ → A\n"
         "- 社内のCcメール → C\n"
         "- ニュースレター・通知 → D\n\n"
         "最後に「今日まずやるべき3つ」を箇条書きで。"),
    ],
)

# 4.4 Gmail応用：返信ドラフト作成
make_single_prompts_slide(
    day_num=4,
    title="Gmail応用：返信ドラフト作成",
    subtitle="読むだけでなく、書く準備までやらせる",
    prompts=[
        ("① A判定メールに返信ドラフト",
         "さきほどのリストの A判定メール について、それぞれ私の文体で返信ドラフトを作ってください。\n"
         "ドラフトはこのチャット内に表示するだけで、絶対に送信はしないでください。\n"
         "私が確認して自分でコピペして送ります。"),
        ("② 関係者の傾向分析",
         "過去30日でやり取りが多かった人TOP10と、それぞれの主な相談内容を整理して。"
         "「最近返信が滞っている人」も指摘して。"),
    ],
    warn={"title": "大事な原則",
          "body": "Claude に自動でメール送信させるのは、慣れるまでやめましょう。「下書きまで作る → 自分で送る」が最も安全。"},
)

# 4.5 ハンズオン②：Calendar 接続
make_single_prompts_slide(
    day_num=4,
    title="ハンズオン②：Calendar 接続",
    subtitle="所要 30分 ／ 週次プランニング",
    steps_first=[
        "Settings → Connectors → Google Calendar の Connect",
        "読み取り権限で接続",
        "下のプロンプトを送信",
    ],
    prompts=[
        ("",
         "来週月曜から金曜までのカレンダー予定を取得して、週次プランを作ってください。\n\n"
         "【出力形式】\n"
         "日別に：\n"
         "- その日の主要予定（時間・タイトル・参加者）\n"
         "- その日に集中作業できる空き時間ブロック（30分以上のもの）\n"
         "- 各会議の前に必要な準備（議題・関連資料）\n"
         "- リスク：詰まりすぎている時間帯、移動時間が足りない予定\n\n"
         "【週全体】\n"
         "- 今週の山場（最重要の会議3つ）と理由\n"
         "- 集中作業に充てられる総時間\n"
         "- 削れそうな会議候補（私が判断するため、削減候補を3つ提案）\n\n"
         "最後に「日曜夜の自分宛てメモ」として、月曜朝に読み返したい一言を。"),
    ],
)

# 4.6 ハンズオン③：Drive 接続 — steps
make_steps_slide(
    day_num=4,
    title="ハンズオン③：Drive 接続",
    subtitle="所要 30分 ／ 資料横断検索 ／ 次のスライドにプロンプト（職種別）",
    steps=[
        "Settings → Connectors → Google Drive の Connect",
        "権限：「Claudeが指定したファイルのみ読み取り」が最も安全",
        "次のスライドのプロンプトを送信（職種別6種）",
    ],
)

# 4.6b drive_search multi
make_multi_prompt_short_slide(
    day_num=4,
    title="コピペ：Drive 横断検索プロンプト",
    subtitle="あなたの職種版をコピー（検索キーワードは職種別に最適化済み）",
    prompt_key="drive_search",
)

# 4.7 3つの統合プロンプト (Lv2)
make_single_prompts_slide(
    day_num=4,
    title="3つの統合プロンプト（Lv2）",
    subtitle="所要 15分 ／ Gmail × Calendar × Drive を1コマンドで",
    prompts=[
        ("明日の完全ブリーフィング",
         "明日の予定を「会議ごとの完全ブリーフィング」にしてください。\n\n"
         "各会議について：\n"
         "1. カレンダーから時間・場所・参加者を取得\n"
         "2. Gmailから、参加者と過去2週間にやり取りしたメールの要点\n"
         "3. Driveから、議題に関連する最新資料3つ\n"
         "4. 「私がこの会議で達成すべきこと」3つの仮説\n"
         "5. 「私が聞かれそうな質問」TOP3と回答案\n\n"
         "会議ごとに見開きで読めるよう、見出しと余白を整えて出力。"),
    ],
    ok={"title": "体感ポイント",
        "body": "朝の準備15〜30分が、コーヒー1杯の間に終わる。「準備不足の会議」が消える。"},
)

# 4.8 プライバシーと運用ルール
make_cards_grid_slide(
    day_num=4,
    title="プライバシーと運用ルール",
    subtitle="便利と引き換えに気をつけること",
    cards=[
        {"title": "⚠ 必ず守る",
         "body": ["顧客の個人情報をプロンプトに直接書かない",
                  "送信操作は必ず自分の手でクリック",
                  "Claudeに「メールを送って」と頼まない（下書きまで）",
                  "会社の機密度の高いフォルダはConnect範囲から外す"]},
        {"title": "✓ 推奨",
         "body": ["業務用と個人用のClaudeアカウントを分ける",
                  "月1で接続済みConnectorsを棚卸し",
                  "退職・異動時は速やかに切断",
                  "不要になったConnectorは即 Disconnect"]},
    ],
)

# 4.9 本日の宿題
make_checklist_slide(
    day_num=4,
    title="本日の宿題",
    subtitle="明日の Day 5（業務に組み込む）に備える",
    items=[
        "明日朝の業務開始時に、統合プロンプト（明日のブリーフィング）を実際に走らせる",
        "Gmailの未読メール処理を、明日1日Claude経由で完結させてみる",
        "同僚に「Connector使ってみてどうだった？」を共有する話題を1つ作る",
        "明日のために、自分の「定型業務リスト」（毎週／毎月やる作業）を5個ほどメモ",
    ],
    ok={"title": "今日の達成",
        "body": "・3つのコネクタを接続して連携  ・朝のブリーフィングを自動化  ・受信箱の処理を効率化  ・プライバシー運用ルールを把握"},
    link="Day 5 へ進む  →",
)


# -----------------------------------------------------------------------------
# DAY 5 — 自分の業務に組み込む
# -----------------------------------------------------------------------------
slide_day_divider(5)

# 5.1 業務棚卸し
make_single_prompts_slide(
    day_num=5,
    title="① 業務棚卸し",
    subtitle="所要 30分 ／ 自分の仕事を「見える化」する",
    steps_first=[
        "紙 or テキストエディタに、過去2週間で実際にやった仕事を全部書き出す（最低20項目）",
        "各項目に頻度（毎日/毎週/毎月/不定期）と所要時間を付ける",
        "下のプロンプトに貼って、Claude に分類してもらう",
    ],
    prompts=[
        ("",
         "以下は私が過去2週間に行った業務リストです。Claude化のしやすさで分類してください。\n\n"
         "【業務リスト】\n"
         "（ここに自分のリストを貼り付け）\n\n"
         "【分類軸】\n"
         "A：今日からClaudeで完全代替可（読む・要約・転記・定型作成）\n"
         "B：Claudeと協働して時短できる（議事録→ToDo、メール下書き等）\n"
         "C：人の判断が必要だが下準備をClaudeに任せられる\n"
         "D：Claudeでは難しい（対面交渉、機密度が高い等）\n\n"
         "【出力】\n"
         "- 業務ごとにA/B/C/D判定と理由\n"
         "- A判定の業務に、それぞれ使う機能（Project/Skill/Connector）と推奨プロンプト概要\n"
         "- 「最初に着手すべきTOP3」をROI順に推薦"),
    ],
)

# 5.2 1業務をフルClaude化 (LONG multi → 2 slides)
make_multi_prompt_long_slide(
    day_num=5,
    title="② 1業務をフルClaude化（Lv2）",
    subtitle="所要 30分 ／ TOP3の中から1つ選んで完成形まで（職種別テンプレ）",
    prompt_key="full_claudify",
    group="a",
)
make_multi_prompt_long_slide(
    day_num=5,
    title="② 1業務をフルClaude化（続き）",
    subtitle="所要 30分 ／ TOP3の中から1つ選んで完成形まで（職種別テンプレ）",
    prompt_key="full_claudify",
    group="b",
)

# 5.3 同僚向けデモ準備
make_single_prompts_slide(
    day_num=5,
    title="③ 同僚向けデモ準備",
    subtitle="所要 30分 ／ 「自分が学んだ」を「組織の資産」に  ／  3分デモのシナリオを作る",
    prompts=[
        ("",
         "同僚（非エンジニア・Claude未経験者）3〜5人を集めた3分デモのシナリオを設計してください。\n\n"
         "【条件】\n"
         "- 私が実際に画面を共有しながら見せる\n"
         "- 今日設計した「（業務名）」のClaude化を題材に\n"
         "- ビフォー（手作業）→ アフター（Claude化）の対比を3分で\n\n"
         "【出力】\n"
         "1. シナリオ進行表（30秒刻み）：何を話す・何を画面で見せる\n"
         "2. 開始トーク（30秒）：相手の関心を掴む一言\n"
         "3. クライマックス：一番「おお」と言われそうな瞬間\n"
         "4. 締め（30秒）：相手にしてほしいアクション（やってみて等）\n"
         "5. 想定質問TOP5と回答\n"
         "6. デモが失敗した時のリカバリ文\n\n"
         "口頭で読み上げる原稿として、丁寧体で書いて。"),
    ],
)

# 5.4 30日ロードマップ
make_single_prompts_slide(
    day_num=5,
    title="④ 30日ロードマップ",
    subtitle="所要 20分 ／ 学習を定着させる仕掛け",
    prompts=[
        ("",
         "今日からの30日間で、Claudeを業務に定着させるロードマップを作ってください。\n\n"
         "【条件】\n"
         "- 平日のみ・1日30分以内の追加投資で済むこと\n"
         "- 段階的に難易度を上げる\n"
         "- 週末は使わない（疲れる前提）\n\n"
         "【週ごとに含めてほしい要素】\n"
         "- 今週のメイン目標（1つだけ）\n"
         "- 月〜金の日替わりタスク（5つ）\n"
         "- 振り返り質問（金曜に自問する3つ）\n"
         "- 達成バッジ（できたら自分にご褒美的な何か）\n\n"
         "【4週目の最後】\n"
         "- この30日を振り返るプロンプトを1つ\n"
         "- 次の30日に進むかどうかの判断基準\n\n"
         "【トーン】\n"
         "ストイックすぎず、続けやすい雰囲気で。"),
    ],
    ok={"title": "仕上げ",
        "body": "このロードマップをカレンダーに転記。1日目から実行開始。"},
)

# 5.5 継続のコツ
make_cards_grid_slide(
    day_num=5,
    title="継続のコツ：3つだけ",
    subtitle="挫折しないために",
    cards=[
        {"title": "① 1日1プロンプト",
         "body": "「毎日30分使う」より「毎日1回は使う」。連続日数を切らないことが最重要。3日空くと戻れない。"},
        {"title": "② 失敗を保存",
         "body": "うまくいかなかった依頼を別 Project に溜める。月1で見直すと、自分の苦手パターンが見える。"},
        {"title": "③ 仲間を1人",
         "body": "同じく学んでいる同僚を最低1人作る。週1回プロンプトを交換する習慣で2人とも上達速度が倍になる。"},
    ],
    cols=3,
    warn={"title": "避けたいパターン",
          "body": "・完璧な Project を作ろうとして1週間使わない  ・1度のミスで「やっぱり使えない」と判断  ・機能を全部使おうとして混乱（必要な順に1つずつ）"},
)

# 5.6 修了
make_checklist_slide(
    day_num=5,
    title="修了 — 5日間で身につけたもの",
    subtitle="確認チェックリスト",
    items=[
        "自分仕様の Project を最低1つ運用している",
        "Memory に自分の前提情報を覚えさせている",
        "PDFや議事録を「自分の視点で」要約させられる",
        "Artifacts で1枚資料を作って編集できる",
        "/xlsx と /pptx を呼び出して使える",
        "最低1つの Connector を業務で活用している",
        "自分の業務を1つフル Claude 化する設計ができた",
        "同僚に3分デモができる",
        "30日ロードマップを持っている",
    ],
    ok={"title": "修了メッセージ",
        "body": "明日からの仕事が、今日までと違うことを願っています。困ったときは、いちばん最初に作った Project に「困っている」と話しかけることから始めてください。"},
)

# 5.7 万能プロンプト集
make_single_prompts_slide(
    day_num=5,
    title="付録：今すぐ使える万能プロンプト集",
    subtitle="明日からのお守り — 4つの場面別プロンプト",
    prompts=[
        ("朝イチ用",
         "今日のカレンダーとメールを見て、私が朝イチで意識すべき3つを教えて。優先順位の根拠も。"),
        ("会議直前用",
         "これから「（会議名）」がある。参加者と議題を踏まえて、5分で頭に入れるブリーフィングを作って。"),
        ("退社前用",
         "今日やったこと（メール送信履歴とカレンダーから）を整理して、明日朝イチで再開しやすいメモを作って。"),
        ("困った時用",
         "私はいま「（状況）」で困っている。私の立場と過去のやり取りを踏まえて、選択肢を3つ提案して。それぞれの長所短所も。"),
    ],
)


# Closing
slide_closing()

# Finalize footers (page n / total)
finalize_footers()

out = "Claude_Chat_5日間研修_配布用ハンズオン資料.pptx"
prs.save(out)
print(f"saved: {out}  ({len(prs.slides)} slides)")
print(f"day totals: {_day_totals}")
