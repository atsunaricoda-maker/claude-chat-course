# -*- coding: utf-8 -*-
"""
配布用スライド資料 (Google Slides 取り込み用 .pptx)
Claude Chat 5日間ハンズオン研修

色パレット (既存教材のブランド準拠):
  primary       #6C47FF
  primary_dark  #4A2BD4
  accent        #FF6B35
  bg            #F6F5FA
  text          #1A1A2E
  muted         #6B7280
  success       #16A34A
  card_bg       #F8F7FC
  warn_bg       #FFF4EE
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor
from lxml import etree

# --- color palette ------------------------------------------------------------
PRIMARY        = RGBColor(0x6C, 0x47, 0xFF)
PRIMARY_DARK   = RGBColor(0x4A, 0x2B, 0xD4)
ACCENT         = RGBColor(0xFF, 0x6B, 0x35)
BG             = RGBColor(0xF6, 0xF5, 0xFA)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
TEXT           = RGBColor(0x1A, 0x1A, 0x2E)
MUTED          = RGBColor(0x6B, 0x72, 0x80)
SUCCESS        = RGBColor(0x16, 0xA3, 0x4A)
CARD_BG        = RGBColor(0xF8, 0xF7, 0xFC)
WARN_BG        = RGBColor(0xFF, 0xF4, 0xEE)
BORDER         = RGBColor(0xE5, 0xE5, 0xEA)

HEAD_FONT = "Yu Gothic"
BODY_FONT = "Yu Gothic"

# --- presentation setup -------------------------------------------------------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
SW_IN = 13.333
SH_IN = 7.5

BLANK = prs.slide_layouts[6]


# --- helpers ------------------------------------------------------------------
def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _kill_shadow(shape):
    sp = shape._element.spPr
    for el in sp.findall(qn("a:effectLst")):
        sp.remove(el)
    etree.SubElement(sp, qn("a:effectLst"))


def add_rect(slide, x, y, w, h, fill=None, line=None, shadow=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if line is not None:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if not shadow:
        _kill_shadow(shape)
    return shape


def add_round_rect(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if line is not None:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    _kill_shadow(shape)
    return shape


def add_oval(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if line is not None:
        shape.line.color.rgb = line
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    _kill_shadow(shape)
    return shape


def add_text(slide, x, y, w, h, text, *,
             size=14, bold=False, italic=False, color=TEXT,
             align="left", valign="top", font=BODY_FONT,
             char_spacing=None, line_spacing=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.TOP

    lines = text.split("\n") if isinstance(text, str) else text
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT
        if line_spacing:
            p.line_spacing = line_spacing
        if isinstance(ln, str):
            run = p.add_run()
            run.text = ln
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color
            if char_spacing:
                rPr = run._r.get_or_add_rPr()
                rPr.set("spc", str(char_spacing))
        else:
            # ln is a list of (text, opts) for rich-run paragraphs
            for txt, opts in ln:
                run = p.add_run()
                run.text = txt
                run.font.name = opts.get("font", font)
                run.font.size = Pt(opts.get("size", size))
                run.font.bold = opts.get("bold", bold)
                run.font.italic = opts.get("italic", italic)
                run.font.color.rgb = opts.get("color", color)
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=14, color=TEXT,
                bullet_color=None, font=BODY_FONT, line_spacing=1.35):
    """items: list[str] — render as left-aligned bullet list with custom marker."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    bc = bullet_color if bullet_color is not None else PRIMARY
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # marker
        mk = p.add_run()
        mk.text = "▸  "
        mk.font.name = font
        mk.font.size = Pt(size)
        mk.font.bold = True
        mk.font.color.rgb = bc
        # text (supports inline bold via **…**)
        parts = item.split("**")
        for j, part in enumerate(parts):
            if part == "":
                continue
            run = p.add_run()
            run.text = part
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = (j % 2 == 1)
            run.font.color.rgb = color
    return tb


def add_header_corner(slide, label, color=PRIMARY):
    """Small page label in top-right corner — replaces 'accent line under title'."""
    add_text(slide, SW_IN - 2.4, 0.4, 2.0, 0.35, label,
             size=10, color=color, bold=True, align="right",
             font=HEAD_FONT, char_spacing=300)


def add_footer(slide, page_no=None, total=None):
    add_text(slide, 0.5, SH_IN - 0.4, 7, 0.3,
             "Claude Chat 5日間ハンズオン研修 / 配布用資料",
             size=9, color=MUTED, font=BODY_FONT)
    if page_no is not None and total is not None:
        add_text(slide, SW_IN - 1.5, SH_IN - 0.4, 1.0, 0.3,
                 f"{page_no} / {total}",
                 size=9, color=MUTED, font=BODY_FONT, align="right")


# =============================================================================
# slide builders
# =============================================================================

TOTAL_PAGES = 16  # update if slides added/removed


def slide_title():
    s = add_slide()
    set_bg(s, RGBColor(0x1A, 0x14, 0x3D))  # deep purple
    # decorative gradient "stripe" via overlapping ovals (no full-width bar)
    add_oval(s, -2.0, -2.5, 7.0, 7.0, RGBColor(0x6C, 0x47, 0xFF))
    add_oval(s, SW_IN - 4.5, SH_IN - 5.5, 8.0, 8.0, RGBColor(0xFF, 0x6B, 0x35))
    # darken overlay so text reads
    overlay = add_rect(s, 0, 0, SW_IN, SH_IN, fill=RGBColor(0x1A, 0x14, 0x3D))
    overlay.fill.fore_color.rgb = RGBColor(0x1A, 0x14, 0x3D)
    overlay.fill.transparency = 0  # python-pptx workaround: set alpha via XML
    # set fill transparency
    sp = overlay.fill.fore_color._xFill
    # We instead just keep one opaque oval for visual variety; remove overlay by
    # actually making it semi-transparent via XML
    a_srgb = sp.find(qn("a:srgbClr"))
    if a_srgb is not None:
        # add alpha 70%
        alpha = etree.SubElement(a_srgb, qn("a:alpha"))
        alpha.set("val", "75000")  # 75% opacity

    # tag pill
    pill = add_round_rect(s, 1.0, 1.05, 4.4, 0.5, fill=RGBColor(0x2A, 0x20, 0x5C))
    pill.line.color.rgb = RGBColor(0x6C, 0x47, 0xFF)
    pill.line.width = Pt(0.75)
    add_text(s, 1.0, 1.05, 4.4, 0.5,
             "DISTRIBUTION DECK / 配布用スライド",
             size=11, bold=True, color=WHITE, align="center", valign="middle",
             font=HEAD_FONT, char_spacing=300)

    # main title
    add_text(s, 1.0, 1.95, 11.5, 1.4,
             "Claude Chat",
             size=72, bold=True, color=WHITE, font=HEAD_FONT)
    add_text(s, 1.0, 3.05, 11.5, 1.0,
             "5日間ハンズオン研修",
             size=44, bold=True, color=RGBColor(0xFF, 0xB7, 0x96), font=HEAD_FONT)

    # divider
    add_rect(s, 1.0, 4.3, 0.8, 0.06, fill=ACCENT)

    # subtitle
    add_text(s, 1.0, 4.5, 11.5, 0.6,
             "毎日2時間 × 5日で、Claude を「明日からの仕事の相棒」に変える実践カリキュラム",
             size=18, color=WHITE, font=BODY_FONT)

    # meta row
    add_text(s, 1.0, 5.45, 11.5, 0.4,
             "全5日 × 各2時間 ／ 合計10時間 ／ 非エンジニア向け ／ ハンズオン90%",
             size=14, color=RGBColor(0xCA, 0xC1, 0xE8), font=BODY_FONT)

    # bottom bar with course URL
    add_text(s, 1.0, SH_IN - 1.0, 11.5, 0.3,
             "教材URL  https://atsunaricoda-maker.github.io/claude-chat-course/",
             size=11, color=RGBColor(0xCA, 0xC1, 0xE8), font=BODY_FONT)
    add_text(s, 1.0, SH_IN - 0.65, 11.5, 0.3,
             "提案者：（記入欄）  ／  提出日：2026年5月",
             size=11, color=RGBColor(0xCA, 0xC1, 0xE8), font=BODY_FONT)


def slide_agenda():
    s = add_slide()
    set_bg(s, WHITE)
    add_header_corner(s, "AGENDA / 目次")
    add_text(s, 0.7, 0.55, 8, 0.7, "本日のアジェンダ",
             size=36, bold=True, color=TEXT, font=HEAD_FONT)
    add_text(s, 0.7, 1.25, 8, 0.4,
             "5日間カリキュラムのねらい・効果・導入条件をまとめてご説明します。",
             size=14, color=MUTED, font=BODY_FONT)

    items = [
        ("01", "なぜ今 Claude か",       "業務における生成AI活用の現状と「研修が業務に活きない」課題"),
        ("02", "研修の全体像",            "5日間カリキュラム構成と特徴"),
        ("03", "各 Day の到達目標",       "Day1〜Day5 で「できる」になる業務"),
        ("04", "他研修との違い",          "ハンズオン90% × 6職種別プロンプト × 30日定着"),
        ("05", "期待効果と ROI",          "1人あたり月21時間削減 / 投資1ヶ月で回収"),
        ("06", "導入要件・スケジュール",  "実施フロー、必要なもの、リスクと対策"),
    ]
    base_y = 1.85
    row_h = 0.78
    for i, (num, title, desc) in enumerate(items):
        y = base_y + i * row_h
        # number circle
        add_oval(s, 0.7, y + 0.05, 0.55, 0.55, PRIMARY)
        add_text(s, 0.7, y + 0.05, 0.55, 0.55, num,
                 size=15, bold=True, color=WHITE, align="center", valign="middle",
                 font=HEAD_FONT)
        # title
        add_text(s, 1.5, y, 4.5, 0.35, title,
                 size=18, bold=True, color=TEXT, font=HEAD_FONT)
        # description
        add_text(s, 1.5, y + 0.36, 11.5, 0.35, desc,
                 size=12, color=MUTED, font=BODY_FONT)

    add_footer(s, 2, TOTAL_PAGES)


def slide_why_now():
    s = add_slide()
    set_bg(s, WHITE)
    add_header_corner(s, "01  なぜ今 Claude か")
    add_text(s, 0.7, 0.55, 12, 0.7, "「研修受けた」≠「業務に活きている」",
             size=34, bold=True, color=TEXT, font=HEAD_FONT)
    add_text(s, 0.7, 1.25, 12, 0.4,
             "生成AI活用は当たり前に。しかし現場には3つの壁が残る。",
             size=14, color=MUTED, font=BODY_FONT)

    # 3 problem cards
    cards = [
        ("受講者の半数以上が業務未活用",
         "「研修は受けたが結局使えていない」",
         "一般的な生成AI研修の主な失敗点。"),
        ("ITリテラシーの個人差",
         "ハードルが高く全社推進が困難",
         "ITが苦手な層ほど初動でつまずく。"),
        ("職種別ニーズへの対応不足",
         "営業・人事・経理で使い方は別物",
         "画一的な研修では現場に届かない。"),
    ]
    cw = 3.95
    gap = 0.25
    base_x = 0.7
    base_y = 2.0
    ch = 3.4
    for i, (title, sub, body) in enumerate(cards):
        x = base_x + i * (cw + gap)
        add_round_rect(s, x, base_y, cw, ch, fill=CARD_BG)
        # left bar
        add_rect(s, x, base_y, 0.08, ch, fill=ACCENT)
        # icon area number
        add_oval(s, x + 0.4, base_y + 0.4, 0.7, 0.7, ACCENT)
        add_text(s, x + 0.4, base_y + 0.4, 0.7, 0.7, f"0{i+1}",
                 size=16, bold=True, color=WHITE, align="center", valign="middle",
                 font=HEAD_FONT)
        add_text(s, x + 0.4, base_y + 1.2, cw - 0.6, 0.7, title,
                 size=15, bold=True, color=TEXT, font=HEAD_FONT,
                 line_spacing=1.25)
        add_text(s, x + 0.4, base_y + 1.95, cw - 0.6, 0.5, sub,
                 size=12, bold=True, color=PRIMARY_DARK, font=BODY_FONT)
        add_text(s, x + 0.4, base_y + 2.45, cw - 0.6, 0.85, body,
                 size=11, color=MUTED, font=BODY_FONT, line_spacing=1.4)

    # bottom callout
    add_round_rect(s, 0.7, 5.6, 12.0, 1.3, fill=RGBColor(0xF0, 0xEE, 0xFC))
    add_rect(s, 0.7, 5.6, 0.08, 1.3, fill=PRIMARY)
    add_text(s, 1.05, 5.75, 11.5, 0.4,
             "本研修の意義",
             size=13, bold=True, color=PRIMARY_DARK, font=HEAD_FONT)
    add_text(s, 1.05, 6.15, 11.5, 0.7,
             "ハンズオン中心の設計と6職種別教材で、「使える状態」まで確実に到達。\n30日のフォローロードマップで業務定着まで伴走します。",
             size=13, color=TEXT, font=BODY_FONT, line_spacing=1.45)

    add_footer(s, 3, TOTAL_PAGES)


def slide_overview():
    s = add_slide()
    set_bg(s, WHITE)
    add_header_corner(s, "02  研修の全体像")
    add_text(s, 0.7, 0.55, 12, 0.7, "5日間で段階的に「自走できる状態」へ",
             size=32, bold=True, color=TEXT, font=HEAD_FONT)
    add_text(s, 0.7, 1.25, 12, 0.4,
             "全5日 / 各2時間 / 合計10時間。各回が独立しつつ、Day1の Project / Memory 設定が後続の前提。",
             size=14, color=MUTED, font=BODY_FONT)

    # 5 day cards in a row
    days = [
        ("DAY 1", "自分仕様にする", "Projects / Memory", "毎回説明しなくていい相棒に"),
        ("DAY 2", "読む・まとめる", "Files / Artifacts",  "資料処理を 1/10 の時間に"),
        ("DAY 3", "データ・スライド自動化", "Skill (/xlsx, /pptx)", "Excel・PPT職人を卒業"),
        ("DAY 4", "ツール連携", "Connector (Gmail等)", "アプリ往復をなくす"),
        ("DAY 5", "業務に組み込む", "棚卸し / 30日設計",   "明日からの定常業務へ"),
    ]
    cw = 2.40
    gap = 0.18
    base_x = 0.7
    base_y = 2.05
    ch = 3.6
    for i, (tag, title, mech, hook) in enumerate(days):
        x = base_x + i * (cw + gap)
        # main card
        add_rect(s, x, base_y, cw, ch, fill=WHITE, line=BORDER)
        # top color block
        col = [PRIMARY, RGBColor(0x84, 0x55, 0xFF), RGBColor(0xB6, 0x52, 0xC9), RGBColor(0xE5, 0x5F, 0x8F), ACCENT][i]
        add_rect(s, x, base_y, cw, 0.7, fill=col)
        add_text(s, x, base_y + 0.13, cw, 0.45, tag,
                 size=14, bold=True, color=WHITE, align="center", font=HEAD_FONT,
                 char_spacing=300)
        # title
        add_text(s, x + 0.2, base_y + 0.85, cw - 0.4, 0.85, title,
                 size=18, bold=True, color=TEXT, font=HEAD_FONT, line_spacing=1.2)
        # mech (the feature used)
        add_text(s, x + 0.2, base_y + 1.85, cw - 0.4, 0.4,
                 "Feature",
                 size=9, bold=True, color=MUTED, font=HEAD_FONT, char_spacing=200)
        add_text(s, x + 0.2, base_y + 2.15, cw - 0.4, 0.6, mech,
                 size=12, bold=True, color=PRIMARY_DARK, font=BODY_FONT, line_spacing=1.3)
        # bottom hook
        add_rect(s, x + 0.2, base_y + 2.85, cw - 0.4, 0.04, fill=BORDER)
        add_text(s, x + 0.2, base_y + 2.95, cw - 0.4, 0.6, hook,
                 size=11, color=MUTED, font=BODY_FONT, line_spacing=1.35)

    # bottom strip — what you get at the end
    add_round_rect(s, 0.7, 6.0, 12.0, 0.95, fill=CARD_BG)
    add_text(s, 0.95, 6.15, 11.5, 0.35, "5日後のあなた",
             size=12, bold=True, color=PRIMARY_DARK, font=HEAD_FONT)
    add_text(s, 0.95, 6.5, 11.5, 0.4,
             "自分の業務文脈を理解した「Claude相棒」を持ち、月21時間の業務時間を創出する状態。",
             size=13, color=TEXT, font=BODY_FONT)

    add_footer(s, 4, TOTAL_PAGES)


def day_detail_slide(day_num, day_title, subtitle, lead, items, colorL, colorR, accent_color, page):
    """共通のDay詳細スライド (左カラム=ねらい/学ぶ機能、右カラム=ハンズオン明細)"""
    s = add_slide()
    set_bg(s, WHITE)
    add_header_corner(s, f"03  各Dayの到達目標", color=accent_color)

    # large day tag
    add_text(s, 0.7, 0.5, 3.0, 0.4, f"DAY  {day_num}",
             size=14, bold=True, color=accent_color, font=HEAD_FONT, char_spacing=400)
    add_text(s, 0.7, 0.85, 12, 1.0, day_title,
             size=32, bold=True, color=TEXT, font=HEAD_FONT, line_spacing=1.15)
    add_text(s, 0.7, 1.85, 12, 0.4, subtitle,
             size=14, color=MUTED, font=BODY_FONT)

    # left column — what you get
    add_round_rect(s, 0.7, 2.55, 5.6, 4.05, fill=CARD_BG)
    add_rect(s, 0.7, 2.55, 0.08, 4.05, fill=accent_color)
    add_text(s, 0.95, 2.7, 5.2, 0.4, "ねらい",
             size=12, bold=True, color=PRIMARY_DARK, font=HEAD_FONT, char_spacing=200)
    add_text(s, 0.95, 3.05, 5.2, 1.7, lead,
             size=14, color=TEXT, font=BODY_FONT, line_spacing=1.55)

    # learning artifacts cards (mini)
    add_text(s, 0.95, 4.85, 5.2, 0.4, "受講後にできること",
             size=12, bold=True, color=PRIMARY_DARK, font=HEAD_FONT, char_spacing=200)
    out_items = items.get("outcomes", [])
    for j, it in enumerate(out_items):
        add_text(s, 0.95, 5.2 + j * 0.32, 5.2, 0.3, f"✓  {it}",
                 size=12, color=TEXT, font=BODY_FONT)

    # right column — agenda
    add_text(s, 6.55, 2.55, 6.2, 0.4, "本編アジェンダ (2時間)",
             size=12, bold=True, color=PRIMARY_DARK, font=HEAD_FONT, char_spacing=200)
    rows = items["agenda"]
    by = 2.95
    rh = 0.62
    for i, (mins, label, hint) in enumerate(rows):
        y = by + i * rh
        # time pill
        pill = add_round_rect(s, 6.55, y, 0.95, 0.45, fill=accent_color)
        add_text(s, 6.55, y, 0.95, 0.45, mins,
                 size=11, bold=True, color=WHITE, align="center", valign="middle",
                 font=HEAD_FONT)
        # label + hint
        add_text(s, 7.65, y - 0.02, 5.1, 0.32, label,
                 size=13, bold=True, color=TEXT, font=HEAD_FONT)
        add_text(s, 7.65, y + 0.28, 5.1, 0.3, hint,
                 size=10, color=MUTED, font=BODY_FONT)

    add_footer(s, page, TOTAL_PAGES)


def slide_day1():
    day_detail_slide(
        day_num="1",
        day_title="Claude を自分仕様にする",
        subtitle="Projects・Memory で「毎回説明しなくていい相棒」を作る。",
        lead=(
            "Project機能・Memory機能を使い、Claudeにあなたの業務文脈や文体を学習させます。"
            "翌日以降のすべての演習が「自分仕様のClaude」上で回るための、最初の土台です。"
        ),
        items={
            "outcomes": [
                "自分専用Projectが立ち上がっている",
                "過去文書から文体を学習させた状態",
                "Memoryに業務基本情報を登録済み",
            ],
            "agenda": [
                ("15分", "アカウント・UI 案内",       "Claude.ai のレイアウトをまず把握"),
                ("30分", "Project 作成 / 指示文設定",  "業務コンテキストをセット"),
                ("35分", "自分の文体を学習させる",     "過去メール・議事録を投入"),
                ("25分", "Memory で個人情報を記憶",    "肩書・部署・取引先など"),
                ("15分", "振り返り・宿題",            "翌日までに準備するもの"),
            ],
        },
        colorL=PRIMARY, colorR=PRIMARY, accent_color=PRIMARY, page=5,
    )


def slide_day2():
    day_detail_slide(
        day_num="2",
        day_title="「読む」「まとめる」を委任する",
        subtitle="Files・Artifacts で資料処理を 1/10 の時間に。",
        lead=(
            "PDF・議事録・録音テキストを Claude に渡し、要約・ToDo化・1枚資料化を一気通貫で。"
            "「自分視点での要約」を引き出す問いの作り方が、この日のコアスキルです。"
        ),
        items={
            "outcomes": [
                "PDFを自分視点で要約できる",
                "議事録→ToDo・共有メールへ自動変換",
                "Artifactsで1枚A4サマリーを作れる",
            ],
            "agenda": [
                ("15分", "ファイル投入の基本",          "PDF / 画像 / テキスト"),
                ("25分", "PDFを自分視点で要約",        "立場と質問軸を指定"),
                ("30分", "議事録→ToDo / 共有メール",   "アウトプットを2段で設計"),
                ("30分", "Artifactsで1枚資料",          "編集を「会話」で重ねる"),
                ("20分", "振り返り・宿題",            "実務PDFで翌日までに練習"),
            ],
        },
        colorL=PRIMARY, colorR=PRIMARY, accent_color=RGBColor(0x84, 0x55, 0xFF), page=6,
    )


def slide_day3():
    day_detail_slide(
        day_num="3",
        day_title="データとスライドを自動化する",
        subtitle="Skill (/xlsx, /pptx) で Excel職人・PPT職人を卒業。",
        lead=(
            "Skill = Claudeの「特殊技能」。/xlsx で月次レポート、/pptx でプレゼン作成を自動化。"
            "Skillの概念→設定→ハンズオンを一気通貫で習得します。"
        ),
        items={
            "outcomes": [
                "Skillの仕組みと設定方法を理解",
                "/xlsx で月次レポートを自動生成",
                "/pptx で15枚プレゼンを15分で",
            ],
            "agenda": [
                ("15分", "Skill とは何か",              "概念・カタログ・ベネフィット"),
                ("15分", "Skill の設定方法・FAQ",       "オン/オフと注意点"),
                ("30分", "/xlsx で月次レポート",        "表とグラフを自動化"),
                ("30分", "グラフ・分析の追加",          "示唆出しまで委ねる"),
                ("30分", "/pptx で15枚スライド",        "ブランド準拠デザインも"),
            ],
        },
        colorL=PRIMARY, colorR=PRIMARY, accent_color=RGBColor(0xB6, 0x52, 0xC9), page=7,
    )


def slide_day4():
    day_detail_slide(
        day_num="4",
        day_title="周辺ツールと繋いで「アプリ往復」を消す",
        subtitle="Connector で Gmail・Calendar・Drive を1つの相棒に。",
        lead=(
            "Connector = Claudeを外部サービスへ接続する仕組み。Gmailの未読要約、Calendarの週次プラン、"
            "Driveの横断検索までを Claude 上で完結させます。プライバシー運用ルールも合わせて確認します。"
        ),
        items={
            "outcomes": [
                "Connectorの概念と設定が分かる",
                "Gmail要約・返信ドラフトを実装",
                "3ツール統合プロンプトを使える",
            ],
            "agenda": [
                ("15分", "Connector とは何か",          "概念・カタログ・ベネフィット"),
                ("15分", "接続方法・FAQ",                "プライバシー運用ルール"),
                ("25分", "Gmail 接続と要約",            "未読・返信ドラフト"),
                ("25分", "Calendar で週次プラン",       "翌週ブリーフィング"),
                ("20分", "Drive で資料横断検索",         "関連ファイルの自動収集"),
                ("20分", "3ツール統合プロンプト",        "横断ブリーフィング"),
            ],
        },
        colorL=PRIMARY, colorR=PRIMARY, accent_color=RGBColor(0xE5, 0x5F, 0x8F), page=8,
    )


def slide_day5():
    day_detail_slide(
        day_num="5",
        day_title="自分の業務に組み込む",
        subtitle="学んだ機能を「明日からの定常業務」に変える。",
        lead=(
            "業務棚卸しワーク (A〜D分類) で自分の仕事をClaude化候補に分けます。"
            "1業務をフルClaude化する詳細設計、同僚向け3分デモ、30日定着ロードマップまで作って終了。"
        ),
        items={
            "outcomes": [
                "業務棚卸し表が完成している",
                "1業務のフルClaude化設計が手元にある",
                "30日ロードマップを実行に移せる",
            ],
            "agenda": [
                ("25分", "業務棚卸しワーク (A〜D)",      "Claude化候補を可視化"),
                ("30分", "1業務をフルClaude化",          "プロンプト+運用を設計"),
                ("25分", "同僚向け3分デモ準備",          "組織内推進者を育てる"),
                ("25分", "30日定着ロードマップ",          "週次で自走へ"),
                ("15分", "修了 / 万能プロンプト集",       "現場で困った時のお守り"),
            ],
        },
        colorL=PRIMARY, colorR=PRIMARY, accent_color=ACCENT, page=9,
    )


def slide_diff():
    s = add_slide()
    set_bg(s, WHITE)
    add_header_corner(s, "04  他研修との違い")
    add_text(s, 0.7, 0.55, 12, 0.7, "「使える状態」まで届く設計",
             size=32, bold=True, color=TEXT, font=HEAD_FONT)
    add_text(s, 0.7, 1.25, 12, 0.4,
             "一般的な生成AI研修と、本研修との比較。", size=14, color=MUTED, font=BODY_FONT)

    rows = [
        ("形式",         "講義中心",                      "ハンズオン90%"),
        ("プロンプト",    "一律のサンプル",                 "営業・マーケ・人事・経理・総務・管理職の6職種別"),
        ("期間構成",     "1日完結",                       "5日間で段階的に習得"),
        ("研修後",       "フォローなし",                   "30日ロードマップ付き"),
        ("教材形態",     "紙資料",                        "Web常時アクセス可・スマホ対応"),
        ("再現性",       "講師の力量に依存",                "コピペ可能なプロンプト集で誰でも再実施"),
        ("UI言語対応",   "英語表記のみ",                   "英語UIラベルに日本語併記"),
    ]

    # header
    by = 2.0
    rh = 0.55
    col1_x = 0.7;  col1_w = 2.4
    col2_x = 3.2;  col2_w = 4.5
    col3_x = 7.85; col3_w = 4.85

    # column headers (with backgrounds)
    add_rect(s, col1_x, by, col1_w, rh, fill=CARD_BG)
    add_rect(s, col2_x, by, col2_w, rh, fill=CARD_BG)
    add_rect(s, col3_x, by, col3_w, rh, fill=RGBColor(0xF0, 0xEE, 0xFC))
    add_text(s, col1_x + 0.2, by, col1_w - 0.2, rh, "観点",
             size=12, bold=True, color=PRIMARY_DARK, valign="middle", font=HEAD_FONT)
    add_text(s, col2_x + 0.2, by, col2_w - 0.2, rh, "一般的な研修",
             size=12, bold=True, color=MUTED, valign="middle", font=HEAD_FONT)
    add_text(s, col3_x + 0.2, by, col3_w - 0.2, rh, "本研修",
             size=12, bold=True, color=PRIMARY_DARK, valign="middle", font=HEAD_FONT)

    # body rows
    for i, (k, v1, v2) in enumerate(rows):
        y = by + (i + 1) * rh
        if i % 2 == 0:
            add_rect(s, col1_x, y, col1_w + col2_w + col3_x - (col1_x + col1_w + col2_w) + col3_w + 0, rh,
                     fill=RGBColor(0xFB, 0xFA, 0xFE))
        # divider line
        add_rect(s, col1_x, y + rh, col1_w + col2_w + col3_w + (col2_x - col1_x - col1_w) + (col3_x - col2_x - col2_w),
                 0.005, fill=BORDER)

        add_text(s, col1_x + 0.2, y, col1_w - 0.2, rh, k,
                 size=12, bold=True, color=TEXT, valign="middle", font=BODY_FONT)
        add_text(s, col2_x + 0.2, y, col2_w - 0.2, rh, v1,
                 size=12, color=MUTED, valign="middle", font=BODY_FONT)
        add_text(s, col3_x + 0.2, y, col3_w - 0.2, rh, v2,
                 size=12, bold=True, color=TEXT, valign="middle", font=BODY_FONT)

    add_footer(s, 10, TOTAL_PAGES)


def slide_outcomes():
    s = add_slide()
    set_bg(s, WHITE)
    add_header_corner(s, "05  期待効果 (受講者個人)")
    add_text(s, 0.7, 0.55, 12, 0.7, "1人あたり 月 21時間 を創出",
             size=32, bold=True, color=TEXT, font=HEAD_FONT)
    add_text(s, 0.7, 1.25, 12, 0.4,
             "受講後の業務時間削減試算 (1人 / 月 / 5業務)。",
             size=14, color=MUTED, font=BODY_FONT)

    rows = [
        ("議事録作成",         "90分 × 4回 = 360分",  "10分 × 4回 = 40分",   "320分"),
        ("月次レポート作成",    "240分 × 1回",         "60分 × 1回",          "180分"),
        ("メール返信",         "30分 × 20日 = 600分",  "10分 × 20日 = 200分", "400分"),
        ("提案書作成",         "180分 × 4本 = 720分",  "60分 × 4本 = 240分",  "480分"),
        ("議事録共有メール",    "30分 × 4回",          "5分 × 4回",            "100分"),
    ]

    # column geometry
    by = 2.05
    rh = 0.55
    cw_label = 2.4
    cw_state = 3.5
    cw_after = 3.5
    cw_save  = 1.8
    x0 = 0.7
    x1 = x0 + cw_label
    x2 = x1 + cw_state
    x3 = x2 + cw_after

    # header
    add_rect(s, x0, by, cw_label + cw_state + cw_after + cw_save, rh, fill=PRIMARY_DARK)
    add_text(s, x0 + 0.2, by, cw_label - 0.2, rh, "業務",
             size=12, bold=True, color=WHITE, valign="middle", font=HEAD_FONT)
    add_text(s, x1 + 0.2, by, cw_state - 0.2, rh, "現状",
             size=12, bold=True, color=WHITE, valign="middle", font=HEAD_FONT)
    add_text(s, x2 + 0.2, by, cw_after - 0.2, rh, "Claude化後",
             size=12, bold=True, color=WHITE, valign="middle", font=HEAD_FONT)
    add_text(s, x3 + 0.2, by, cw_save - 0.2, rh, "削減/月",
             size=12, bold=True, color=WHITE, valign="middle", font=HEAD_FONT, align="right")

    # body
    for i, (label, before, after, save) in enumerate(rows):
        y = by + (i + 1) * rh
        if i % 2 == 0:
            add_rect(s, x0, y, cw_label + cw_state + cw_after + cw_save, rh, fill=CARD_BG)
        add_text(s, x0 + 0.2, y, cw_label - 0.2, rh, label,
                 size=12, bold=True, color=TEXT, valign="middle", font=BODY_FONT)
        add_text(s, x1 + 0.2, y, cw_state - 0.2, rh, before,
                 size=11, color=MUTED, valign="middle", font=BODY_FONT)
        add_text(s, x2 + 0.2, y, cw_after - 0.2, rh, after,
                 size=11, color=TEXT, valign="middle", font=BODY_FONT)
        add_text(s, x3 + 0.2, y, cw_save - 0.2, rh, save,
                 size=12, bold=True, color=SUCCESS, valign="middle", font=HEAD_FONT, align="right")

    # total row
    y = by + (len(rows) + 1) * rh
    add_rect(s, x0, y, cw_label + cw_state + cw_after + cw_save, rh + 0.05, fill=RGBColor(0xF0, 0xEE, 0xFC))
    add_text(s, x0 + 0.2, y, cw_label - 0.2, rh, "合計（月）",
             size=13, bold=True, color=PRIMARY_DARK, valign="middle", font=HEAD_FONT)
    add_text(s, x1 + 0.2, y, cw_state - 0.2, rh, "約 33時間",
             size=12, bold=True, color=MUTED, valign="middle", font=BODY_FONT)
    add_text(s, x2 + 0.2, y, cw_after - 0.2, rh, "約 12時間",
             size=12, bold=True, color=TEXT, valign="middle", font=BODY_FONT)
    add_text(s, x3 + 0.2, y, cw_save - 0.2, rh, "約 21時間",
             size=15, bold=True, color=ACCENT, valign="middle", font=HEAD_FONT, align="right")

    # qualitative note (placed below the table, before footer)
    add_round_rect(s, 0.7, 6.10, 12.0, 0.85, fill=WARN_BG)
    add_rect(s, 0.7, 6.10, 0.08, 0.85, fill=ACCENT)
    add_text(s, 1.05, 6.18, 11.5, 0.32, "定性効果",
             size=12, bold=True, color=PRIMARY_DARK, font=HEAD_FONT)
    add_text(s, 1.05, 6.50, 11.5, 0.4,
             "文書品質の標準化 ／ 「白紙が怖い」状態からの脱却 ／ ツール変化への自走力 ／ 同僚への伝授で社内講師化",
             size=11, color=TEXT, font=BODY_FONT)

    add_footer(s, 11, TOTAL_PAGES)


def slide_roi():
    s = add_slide()
    set_bg(s, WHITE)
    add_header_corner(s, "06  ROI 試算")
    add_text(s, 0.7, 0.55, 12, 0.7, "投資1ヶ月で回収する高ROI施策",
             size=32, bold=True, color=TEXT, font=HEAD_FONT)
    add_text(s, 0.7, 1.25, 12, 0.4,
             "受講者30名規模での試算 (規模・条件に応じ調整可能)。",
             size=14, color=MUTED, font=BODY_FONT)

    # two big number cards
    cw = 5.85
    ch = 2.6
    bx = 0.7
    by = 2.0

    # invest
    add_round_rect(s, bx, by, cw, ch, fill=WARN_BG)
    add_rect(s, bx, by, 0.08, ch, fill=ACCENT)
    add_text(s, bx + 0.4, by + 0.3, cw - 0.4, 0.4, "投資 (年間)",
             size=14, bold=True, color=PRIMARY_DARK, font=HEAD_FONT)
    add_text(s, bx + 0.4, by + 0.85, cw - 0.4, 1.0, "約 100 万円",
             size=58, bold=True, color=TEXT, font=HEAD_FONT)
    add_text(s, bx + 0.4, by + 1.95, cw - 0.4, 0.6,
             "講師費 + ツール費 + 機会費用\n教材費は ¥0 (全Web公開済み)",
             size=12, color=MUTED, font=BODY_FONT, line_spacing=1.4)

    # value
    add_round_rect(s, bx + cw + 0.3, by, cw, ch, fill=RGBColor(0xF0, 0xFD, 0xF4))
    add_rect(s, bx + cw + 0.3, by, 0.08, ch, fill=SUCCESS)
    add_text(s, bx + cw + 0.7, by + 0.3, cw - 0.4, 0.4, "創出付加価値 (年間)",
             size=14, bold=True, color=PRIMARY_DARK, font=HEAD_FONT)
    add_text(s, bx + cw + 0.7, by + 0.85, cw - 0.4, 1.0, "約 1,440 万円",
             size=58, bold=True, color=SUCCESS, font=HEAD_FONT)
    add_text(s, bx + cw + 0.7, by + 1.95, cw - 0.4, 0.6,
             "1人あたり年間192時間 × 時給2,500円 × 30名\n= 14,400,000円/年",
             size=12, color=MUTED, font=BODY_FONT, line_spacing=1.4)

    # calc breakdown
    add_text(s, 0.7, 4.80, 12, 0.35, "計算根拠",
             size=13, bold=True, color=PRIMARY_DARK, font=HEAD_FONT)
    bullets = [
        "1人あたり月間時間削減：**21時間** (前項試算)",
        "年間時間削減：21時間 × 12ヶ月 = **252時間/人/年** (保守的に192時間で計算)",
        "時給換算：年収500万円相当として **2,500円/時**",
        "30名規模：480,000円 × 30 = **14,400,000円/年**",
    ]
    add_bullets(s, 0.7, 5.15, 12, 1.25, bullets, size=11, color=TEXT, line_spacing=1.3)

    # callout (placed above footer)
    add_round_rect(s, 0.7, 6.50, 12.0, 0.5, fill=PRIMARY_DARK)
    add_text(s, 0.95, 6.50, 11.5, 0.5,
             "回収期間：約1ヶ月。投資100万円に対し、月間120万円の付加価値創出により、初月でほぼ回収完了。",
             size=12, bold=True, color=WHITE, valign="middle", font=BODY_FONT)

    add_footer(s, 12, TOTAL_PAGES)


def slide_org_impact():
    s = add_slide()
    set_bg(s, WHITE)
    add_header_corner(s, "06  組織・社会への波及")
    add_text(s, 0.7, 0.55, 12, 0.7, "個人の時間創出が、組織の競争力に変わる",
             size=32, bold=True, color=TEXT, font=HEAD_FONT)
    add_text(s, 0.7, 1.25, 12, 0.4,
             "受講者個人の効率化を、組織全体・社会的便益にどう接続するか。",
             size=14, color=MUTED, font=BODY_FONT)

    # 2 columns
    cx = 0.7
    cw = 5.95
    cy = 2.0
    ch = 4.6

    # org column
    add_round_rect(s, cx, cy, cw, ch, fill=CARD_BG)
    add_rect(s, cx, cy, 0.08, ch, fill=PRIMARY)
    add_text(s, cx + 0.35, cy + 0.3, cw - 0.5, 0.5, "組織への波及効果",
             size=18, bold=True, color=PRIMARY_DARK, font=HEAD_FONT)
    org_items = [
        "**月21時間 × 受講者数** = 大きな時間創出",
        "創出時間を市民サービス・企画立案・新規事業へ転用",
        "部署横断のベストプラクティス共有による組織学習",
        "新人育成期間の短縮、属人化解消",
        "残業時間削減・ワークライフバランス改善",
    ]
    add_bullets(s, cx + 0.35, cy + 1.0, cw - 0.5, ch - 1.0, org_items,
                size=13, line_spacing=1.55)

    # social column
    add_round_rect(s, cx + cw + 0.3, cy, cw, ch, fill=WARN_BG)
    add_rect(s, cx + cw + 0.3, cy, 0.08, ch, fill=ACCENT)
    add_text(s, cx + cw + 0.65, cy + 0.3, cw - 0.5, 0.5, "社会的便益",
             size=18, bold=True, color=PRIMARY_DARK, font=HEAD_FONT)
    soc_items = [
        "職員の働き方改革推進",
        "公務効率化による財政負担の軽減",
        "**先進事例**として他自治体・他組織への波及",
        "DX推進指標 (KPI) への直接的貢献",
        "市民への迅速な対応・サービス品質向上",
    ]
    add_bullets(s, cx + cw + 0.65, cy + 1.0, cw - 0.5, ch - 1.0, soc_items,
                size=13, line_spacing=1.55, bullet_color=ACCENT)

    add_footer(s, 13, TOTAL_PAGES)


def slide_implementation():
    s = add_slide()
    set_bg(s, WHITE)
    add_header_corner(s, "07  導入要件")
    add_text(s, 0.7, 0.55, 12, 0.7, "「明日から始められる」軽量な導入要件",
             size=32, bold=True, color=TEXT, font=HEAD_FONT)
    add_text(s, 0.7, 1.25, 12, 0.4,
             "教材は全てWeb公開済み。Claude.aiのアカウントさえあれば即実施可能。",
             size=14, color=MUTED, font=BODY_FONT)

    cards = [
        ("必要なもの", PRIMARY, [
            "受講者用PC (インターネット接続)",
            "Claude.ai アカウント (無料〜Enterprise)",
            "会議室・プロジェクター (集合形式時)",
            "ビデオ会議ツール (オンライン形式時)",
        ]),
        ("事前準備 (推奨)", RGBColor(0x84, 0x55, 0xFF), [
            "受講者の業務カテゴリ事前ヒアリング",
            "機密情報取扱いポリシーのすり合わせ",
            "情シス部門との接続確認 (Connector)",
            "実業務の素材 (メール・議事録) 準備依頼",
        ]),
        ("講師", ACCENT, [
            "本カリキュラム作成者または認定講師",
            "非エンジニア研修の経験者",
            "受講者の質問にリアルタイム対応",
            "30日後のフォローアップを担当",
        ]),
    ]

    cw = 4.05
    gap = 0.18
    bx = 0.7
    by = 2.0
    ch = 4.6
    for i, (title, col, items) in enumerate(cards):
        x = bx + i * (cw + gap)
        add_round_rect(s, x, by, cw, ch, fill=CARD_BG)
        add_rect(s, x, by, 0.08, ch, fill=col)
        add_oval(s, x + 0.35, by + 0.35, 0.55, 0.55, col)
        add_text(s, x + 0.35, by + 0.35, 0.55, 0.55, str(i + 1),
                 size=15, bold=True, color=WHITE, align="center", valign="middle",
                 font=HEAD_FONT)
        add_text(s, x + 1.1, by + 0.4, cw - 1.2, 0.5, title,
                 size=17, bold=True, color=TEXT, font=HEAD_FONT)
        for j, it in enumerate(items):
            yy = by + 1.25 + j * 0.75
            add_text(s, x + 0.4, yy, cw - 0.6, 0.7, f"✓  {it}",
                     size=12, color=TEXT, font=BODY_FONT, line_spacing=1.4)

    add_footer(s, 14, TOTAL_PAGES)


def slide_schedule():
    s = add_slide()
    set_bg(s, WHITE)
    add_header_corner(s, "07  実施スケジュール例")
    add_text(s, 0.7, 0.55, 12, 0.7, "標準は2週間プログラム / 柔軟に調整可",
             size=32, bold=True, color=TEXT, font=HEAD_FONT)
    add_text(s, 0.7, 1.25, 12, 0.4,
             "週2〜3回の頻度で2週間。短期集中・分散型の調整も可能です。",
             size=14, color=MUTED, font=BODY_FONT)

    # timeline
    rows = [
        ("第1週", "月", "Day 1 — Claudeを自分仕様にする",     "2時間", PRIMARY),
        ("第1週", "水", "Day 2 — 「読む」「まとめる」を委任",  "2時間", RGBColor(0x84, 0x55, 0xFF)),
        ("第1週", "金", "Day 3 — データとスライドを自動化",     "2時間", RGBColor(0xB6, 0x52, 0xC9)),
        ("第2週", "火", "Day 4 — 周辺ツールと繋ぐ",             "2時間", RGBColor(0xE5, 0x5F, 0x8F)),
        ("第2週", "木", "Day 5 — 自分の業務に組み込む",          "2時間", ACCENT),
        ("実行期", "—", "受講者各自で30日ロードマップを実行",     "30日", MUTED),
        ("30日後", "—", "フォローアップ (希望者対象)",            "30分", PRIMARY_DARK),
    ]
    by = 2.0
    rh = 0.55
    for i, (week, day, label, dur, col) in enumerate(rows):
        y = by + i * rh
        # week pill
        add_rect(s, 0.7, y, 1.2, rh - 0.05, fill=CARD_BG)
        add_text(s, 0.7, y, 1.2, rh - 0.05, week,
                 size=11, bold=True, color=PRIMARY_DARK, valign="middle", align="center", font=HEAD_FONT)
        # day box
        add_oval(s, 2.0, y + 0.06, rh - 0.15, rh - 0.15, col)
        add_text(s, 2.0, y + 0.06, rh - 0.15, rh - 0.15, day,
                 size=11, bold=True, color=WHITE, align="center", valign="middle", font=HEAD_FONT)
        # label
        add_text(s, 2.7, y, 8.5, rh - 0.05, label,
                 size=13, color=TEXT, valign="middle", font=BODY_FONT)
        # duration
        add_text(s, 11.3, y, 1.4, rh - 0.05, dur,
                 size=12, bold=True, color=col, valign="middle", align="right", font=HEAD_FONT)

    # flexibility
    add_round_rect(s, 0.7, 6.05, 12.0, 0.95, fill=CARD_BG)
    add_rect(s, 0.7, 6.05, 0.08, 0.95, fill=PRIMARY)
    add_text(s, 1.05, 6.15, 11.5, 0.35, "柔軟な調整",
             size=12, bold=True, color=PRIMARY_DARK, font=HEAD_FONT)
    add_text(s, 1.05, 6.5, 11.5, 0.4,
             "1日に2 Day連続 (1週間圧縮) / 1ヶ月に1 Day (ゆったり進行) / 部署別開催 / 全社一括 — いずれも対応可。",
             size=12, color=TEXT, font=BODY_FONT)

    add_footer(s, 15, TOTAL_PAGES)


def slide_closing():
    s = add_slide()
    set_bg(s, RGBColor(0x1A, 0x14, 0x3D))
    # decorative shapes
    add_oval(s, SW_IN - 4.5, -2.5, 7.0, 7.0, RGBColor(0xFF, 0x6B, 0x35))
    add_oval(s, -3.0, SH_IN - 5.0, 8.0, 8.0, RGBColor(0x6C, 0x47, 0xFF))
    overlay = add_rect(s, 0, 0, SW_IN, SH_IN, fill=RGBColor(0x1A, 0x14, 0x3D))
    sp = overlay.fill.fore_color._xFill
    a_srgb = sp.find(qn("a:srgbClr"))
    if a_srgb is not None:
        alpha = etree.SubElement(a_srgb, qn("a:alpha"))
        alpha.set("val", "78000")

    add_text(s, 1.0, 1.7, 11.5, 0.5,
             "CLOSING / 結語",
             size=12, bold=True, color=RGBColor(0xFF, 0xB7, 0x96), font=HEAD_FONT, char_spacing=400)

    add_text(s, 1.0, 2.3, 11.5, 1.4,
             "DX推進・働き方改革・業務効率化",
             size=42, bold=True, color=WHITE, font=HEAD_FONT, line_spacing=1.2)
    add_text(s, 1.0, 3.2, 11.5, 1.0,
             "3つの政策目標に同時に貢献する施策へ",
             size=28, bold=True, color=RGBColor(0xFF, 0xB7, 0x96), font=HEAD_FONT)

    add_text(s, 1.0, 4.4, 11.5, 1.4,
             "5日間×2時間という現実的な時間投資で、職員が自走できる状態に到達。\n"
             "30日のフォローで業務定着まで伴走する仕組みです。",
             size=15, color=RGBColor(0xCA, 0xC1, 0xE8), font=BODY_FONT, line_spacing=1.7)

    # CTA card
    add_round_rect(s, 1.0, 5.85, 11.3, 1.2, fill=WHITE)
    add_text(s, 1.4, 5.95, 10.5, 0.4, "教材は本日からアクセス可能です",
             size=12, bold=True, color=PRIMARY_DARK, font=HEAD_FONT)
    add_text(s, 1.4, 6.3, 10.5, 0.4,
             "https://atsunaricoda-maker.github.io/claude-chat-course/",
             size=18, bold=True, color=PRIMARY_DARK, font=HEAD_FONT)
    add_text(s, 1.4, 6.65, 10.5, 0.35,
             "ご質問・ご相談は提案者まで。提案書 (Gドキュメント) も別添でお送りいたします。",
             size=11, color=MUTED, font=BODY_FONT)


# =============================================================================
# build
# =============================================================================
slide_title()
slide_agenda()
slide_why_now()
slide_overview()
slide_day1()
slide_day2()
slide_day3()
slide_day4()
slide_day5()
slide_diff()
slide_outcomes()
slide_roi()
slide_org_impact()
slide_implementation()
slide_schedule()
slide_closing()

out_path = "Claude_Chat_5日間研修_配布用スライド.pptx"
prs.save(out_path)
print(f"saved: {out_path}  / slides = {len(prs.slides)}")
