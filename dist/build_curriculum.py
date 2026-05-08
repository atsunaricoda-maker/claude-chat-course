# -*- coding: utf-8 -*-
"""
カリキュラム表 (Google Slides 取り込み用 .pptx)
Claude Chat 5日間ハンズオン研修
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor
from lxml import etree

# --- color palette --------------------------------------------------------
PRIMARY        = RGBColor(0x6C, 0x47, 0xFF)
PRIMARY_DARK   = RGBColor(0x4A, 0x2B, 0xD4)
ACCENT         = RGBColor(0xFF, 0x6B, 0x35)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
TEXT           = RGBColor(0x1A, 0x1A, 0x2E)
MUTED          = RGBColor(0x6B, 0x72, 0x80)
SUCCESS        = RGBColor(0x16, 0xA3, 0x4A)
CARD_BG        = RGBColor(0xF8, 0xF7, 0xFC)
WARN_BG        = RGBColor(0xFF, 0xF4, 0xEE)
BORDER         = RGBColor(0xE5, 0xE5, 0xEA)
SUBTLE_GREY    = RGBColor(0xFB, 0xFA, 0xFE)

# Day-specific colors
DAY_COLORS = [
    RGBColor(0x6C, 0x47, 0xFF),  # Day 1 — primary
    RGBColor(0x84, 0x55, 0xFF),  # Day 2
    RGBColor(0xB6, 0x52, 0xC9),  # Day 3
    RGBColor(0xE5, 0x5F, 0x8F),  # Day 4
    RGBColor(0xFF, 0x6B, 0x35),  # Day 5 — accent
]

HEAD_FONT = "Yu Gothic"
BODY_FONT = "Yu Gothic"

# --- presentation setup ---------------------------------------------------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW_IN, SH_IN = 13.333, 7.5
BLANK = prs.slide_layouts[6]


# --- helpers --------------------------------------------------------------
def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _kill_shadow(shape):
    sp = shape._element.spPr
    for el in sp.findall(qn("a:effectLst")):
        sp.remove(el)
    etree.SubElement(sp, qn("a:effectLst"))


def add_rect(slide, x, y, w, h, fill=None, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.line.fill.background()
    if line is not None:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    _kill_shadow(sh)
    return sh


def add_round(slide, x, y, w, h, fill=None, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.line.fill.background()
    if line is not None:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    _kill_shadow(sh)
    return sh


def add_oval(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.line.fill.background()
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    _kill_shadow(sh)
    return sh


def add_text(slide, x, y, w, h, text, *,
             size=14, bold=False, color=TEXT,
             align="left", valign="top", font=BODY_FONT,
             char_spacing=None, line_spacing=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.word_wrap = True
    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.TOP

    lines = text.split("\n") if isinstance(text, str) else text
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = ln
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        if char_spacing:
            rPr = run._r.get_or_add_rPr()
            rPr.set("spc", str(char_spacing))
    return tb


# =========================================================================
# SLIDE 1: マスター時間割（横向き 5列）
# =========================================================================

def slide_master_timetable():
    s = add_slide()
    set_bg(s, WHITE)

    # corner label
    add_text(s, SW_IN - 4.0, 0.4, 3.6, 0.35,
             "CURRICULUM SCHEDULE / カリキュラム表",
             size=10, bold=True, color=PRIMARY, align="right",
             font=HEAD_FONT, char_spacing=300)

    # title
    add_text(s, 0.6, 0.45, 12.0, 0.6,
             "Claude Chat 5日間ハンズオン研修 — 全体時間割",
             size=26, bold=True, color=TEXT, font=HEAD_FONT)
    add_text(s, 0.6, 1.05, 12.0, 0.4,
             "全5日 × 各2時間 ／ 合計10時間 ／ 6職種別プロンプト対応 ／ 30日定着フォロー付",
             size=12, color=MUTED, font=BODY_FONT)

    # Header band — Day labels
    days = [
        ("DAY 1", "Claudeを自分仕様にする",     "Projects・Memory"),
        ("DAY 2", "「読む・まとめる」を委任",   "Files・Artifacts"),
        ("DAY 3", "データ・スライドを自動化",  "Skill (/xlsx,/pptx)"),
        ("DAY 4", "ツール連携で往復をなくす",   "Connector (Gmail等)"),
        ("DAY 5", "自分の業務に組み込む",       "棚卸し・30日設計"),
    ]
    col_x = 0.6
    col_w = 2.46
    col_gap = 0.04
    header_y = 1.7
    header_h = 1.0

    for i, (tag, title, mech) in enumerate(days):
        x = col_x + i * (col_w + col_gap)
        add_rect(s, x, header_y, col_w, header_h, fill=DAY_COLORS[i])
        add_text(s, x, header_y + 0.10, col_w, 0.36, tag,
                 size=14, bold=True, color=WHITE, align="center",
                 font=HEAD_FONT, char_spacing=300)
        add_text(s, x + 0.15, header_y + 0.42, col_w - 0.3, 0.34, title,
                 size=12, bold=True, color=WHITE, align="center", font=HEAD_FONT,
                 line_spacing=1.2)
        add_text(s, x + 0.15, header_y + 0.74, col_w - 0.3, 0.24, mech,
                 size=9, color=RGBColor(0xFF, 0xE7, 0xDD), align="center", font=BODY_FONT)

    # body cells: 5 rows of agendas (the longest day has 6 agenda items, but we'll fit 5 with shading)
    rows = [
        # row format: each cell = (mins, label)
        [
            ("15分", "アカウント・UI 案内"),
            ("15分", "ファイル投入の基本"),
            ("15分", "Skill とは"),
            ("15分", "Connector とは"),
            ("25分", "業務棚卸しワーク"),
        ],
        [
            ("30分", "Project 作成・指示文"),
            ("25分", "PDFを自分視点で要約"),
            ("15分", "Skill 設定方法・FAQ"),
            ("15分", "接続方法・FAQ"),
            ("30分", "1業務をフルClaude化"),
        ],
        [
            ("35分", "自分の文体を学習"),
            ("30分", "議事録 → ToDo / 共有メール"),
            ("30分", "/xlsx で月次レポート"),
            ("25分", "Gmail 接続・要約"),
            ("25分", "同僚向け 3分デモ"),
        ],
        [
            ("25分", "Memory で個人情報記憶"),
            ("30分", "Artifacts で 1枚資料"),
            ("30分", "グラフ・分析の追加"),
            ("25分", "Calendar 週次プラン"),
            ("25分", "30日定着ロードマップ"),
        ],
        [
            ("15分", "振り返り・宿題"),
            ("20分", "宿題・Day3への準備"),
            ("30分", "/pptx で15枚スライド"),
            ("20分", "Drive 横断検索"),
            ("15分", "万能プロンプト集"),
        ],
        [
            ("", ""),  # day1 has no 6th
            ("", ""),  # day2 has no 6th
            ("", ""),  # day3 has no 6th
            ("20分", "3ツール統合プロンプト"),  # day4 has 6th
            ("", ""),  # day5 has no 6th
        ],
    ]

    body_y = 2.85
    row_h = 0.60
    for r_idx, row in enumerate(rows):
        for c_idx, (mins, label) in enumerate(row):
            x = col_x + c_idx * (col_w + col_gap)
            y = body_y + r_idx * row_h
            # background tint alternating
            if r_idx % 2 == 0:
                add_rect(s, x, y, col_w, row_h - 0.05, fill=SUBTLE_GREY)
            else:
                add_rect(s, x, y, col_w, row_h - 0.05, fill=WHITE, line=BORDER)
            if mins:
                # time pill
                add_rect(s, x + 0.1, y + 0.10, 0.55, 0.32, fill=DAY_COLORS[c_idx])
                add_text(s, x + 0.1, y + 0.10, 0.55, 0.32, mins,
                         size=10, bold=True, color=WHITE, align="center", valign="middle",
                         font=HEAD_FONT)
                # label
                add_text(s, x + 0.75, y + 0.06, col_w - 0.85, row_h - 0.15, label,
                         size=11, color=TEXT, valign="middle", font=BODY_FONT,
                         line_spacing=1.25)

    # bottom strip — totals + meta
    total_y = body_y + 6 * row_h + 0.10  # 2.85 + 3.6 + 0.1 = 6.55
    add_round(s, 0.6, total_y, 12.5, 0.55, fill=PRIMARY_DARK)
    add_text(s, 0.85, total_y, 4.0, 0.55, "合計  10時間",
             size=15, bold=True, color=WHITE, valign="middle", font=HEAD_FONT)
    add_text(s, 4.85, total_y, 8.0, 0.55,
             "  ／  各日 2時間  ／  ハンズオン90%  ／  6職種別プロンプト  ／  30日定着フォロー",
             size=11, color=RGBColor(0xCA, 0xC1, 0xE8), valign="middle", font=BODY_FONT)

    # footer
    add_text(s, 0.6, SH_IN - 0.4, 12.5, 0.3,
             "教材URL：https://atsunaricoda-maker.github.io/claude-chat-course/   "
             "／  各セルは2時間枠内の所要時間  ／  6職種：営業 / マーケ・企画 / 人事 / 経理 / 総務・事務 / 管理職",
             size=9, color=MUTED, font=BODY_FONT)


# =========================================================================
# SLIDE 2: Day別の詳細表
# =========================================================================

def slide_day_details():
    s = add_slide()
    set_bg(s, WHITE)

    add_text(s, SW_IN - 4.0, 0.4, 3.6, 0.35,
             "DAY-BY-DAY DETAIL / Day別カリキュラム",
             size=10, bold=True, color=PRIMARY, align="right",
             font=HEAD_FONT, char_spacing=300)
    add_text(s, 0.6, 0.45, 12.0, 0.6,
             "Day別カリキュラム詳細",
             size=26, bold=True, color=TEXT, font=HEAD_FONT)
    add_text(s, 0.6, 1.05, 12.0, 0.4,
             "ねらい・到達目標・主要プロンプト・宿題までを1枚で俯瞰",
             size=12, color=MUTED, font=BODY_FONT)

    # 5 cards horizontally
    days = [
        {
            "tag": "DAY 1",
            "title": "Claudeを自分仕様にする",
            "goal": "Projects・Memoryで「毎回説明しなくていい相棒」を作る",
            "tools": "Projects / Memory",
            "outputs": [
                "自分専用Project運用開始",
                "過去文書から文体学習",
                "Memoryに業務基本情報",
            ],
            "prompts": "・Project指示文（職種別6種）\n・文体学習追記文\n・Memory基本情報（職種別6種）",
            "homework": "メール3〜5通追加投入／PDF・議事録を1つずつ準備",
        },
        {
            "tag": "DAY 2",
            "title": "「読む」「まとめる」を委任",
            "goal": "Files・Artifactsで資料処理を 1/10 の時間に",
            "tools": "ファイル投入 / Artifacts",
            "outputs": [
                "PDFを自分視点で要約",
                "議事録→ToDoリスト変換",
                "Artifactsで1枚A4資料",
            ],
            "prompts": "・PDF要約プロンプト\n・議事録→ToDo→共有メール\n・1枚資料（職種別6種）",
            "homework": "Excel/CSV を1つ準備／PPTテーマを1つ決める",
        },
        {
            "tag": "DAY 3",
            "title": "データ・スライドを自動化",
            "goal": "Skill (/xlsx, /pptx) でExcel・PPT職人を卒業",
            "tools": "Skill (/xlsx, /pptx)",
            "outputs": [
                "Skillの仕組みと設定理解",
                "/xlsx で月次レポート",
                "/pptx で15枚プレゼン",
            ],
            "prompts": "・/xlsx 月次（職種別6種）\n・グラフ・分析プロンプト\n・/pptx 15枚（職種別6種）\n・ブランド準拠プロンプト",
            "homework": "作った成果物を実業務で使う／メール・カレンダー接続確認",
        },
        {
            "tag": "DAY 4",
            "title": "ツール連携で往復をなくす",
            "goal": "Connector で Gmail・Calendar・Drive を1つの相棒に",
            "tools": "Connector (Gmail/Calendar/Drive)",
            "outputs": [
                "Gmail要約・返信ドラフト",
                "Calendar週次プラン",
                "3ツール統合ブリーフィング",
            ],
            "prompts": "・Gmail受信箱要約\n・Gmail返信ドラフト\n・Calendar週次プラン\n・Drive横断検索（職種別6種）\n・3ツール統合プロンプト",
            "homework": "明日朝の統合プロンプト実行／定型業務リスト5個メモ",
        },
        {
            "tag": "DAY 5",
            "title": "自分の業務に組み込む",
            "goal": "学んだ機能を「明日からの定常業務」に変える",
            "tools": "棚卸し / 30日設計",
            "outputs": [
                "業務棚卸し表 (A〜D分類)",
                "1業務のフルClaude化設計",
                "30日定着ロードマップ",
            ],
            "prompts": "・業務棚卸しプロンプト\n・フルClaude化（職種別6種）\n・3分デモシナリオ\n・30日ロードマップ\n・万能プロンプト集 4種",
            "homework": "1日1プロンプト／失敗を別Project／仲間を1人",
        },
    ]

    cw = 2.46
    gap = 0.04
    bx = 0.6
    by = 1.7
    ch = 5.55

    for i, d in enumerate(days):
        x = bx + i * (cw + gap)
        # card outline
        add_rect(s, x, by, cw, ch, fill=WHITE, line=BORDER)
        # header
        add_rect(s, x, by, cw, 0.7, fill=DAY_COLORS[i])
        add_text(s, x, by + 0.10, cw, 0.30, d["tag"],
                 size=12, bold=True, color=WHITE, align="center",
                 font=HEAD_FONT, char_spacing=300)
        add_text(s, x + 0.15, by + 0.40, cw - 0.3, 0.30, d["title"],
                 size=11, bold=True, color=WHITE, align="center", font=HEAD_FONT,
                 line_spacing=1.15)

        # body
        cy = by + 0.85

        # goal
        add_text(s, x + 0.15, cy, cw - 0.3, 0.22, "ねらい",
                 size=8, bold=True, color=PRIMARY_DARK, font=HEAD_FONT, char_spacing=200)
        add_text(s, x + 0.15, cy + 0.22, cw - 0.3, 0.7, d["goal"],
                 size=9.5, color=TEXT, font=BODY_FONT, line_spacing=1.4)

        cy += 0.95

        # tools
        add_text(s, x + 0.15, cy, cw - 0.3, 0.22, "使う機能",
                 size=8, bold=True, color=PRIMARY_DARK, font=HEAD_FONT, char_spacing=200)
        add_text(s, x + 0.15, cy + 0.22, cw - 0.3, 0.32, d["tools"],
                 size=10, bold=True, color=DAY_COLORS[i], font=BODY_FONT)

        cy += 0.6

        # outputs
        add_text(s, x + 0.15, cy, cw - 0.3, 0.22, "受講後にできること",
                 size=8, bold=True, color=PRIMARY_DARK, font=HEAD_FONT, char_spacing=200)
        for j, item in enumerate(d["outputs"]):
            add_text(s, x + 0.15, cy + 0.22 + j * 0.28, cw - 0.3, 0.28,
                     f"✓  {item}",
                     size=9, color=TEXT, font=BODY_FONT)

        cy += 0.22 + 3 * 0.28 + 0.1

        # main prompts
        add_text(s, x + 0.15, cy, cw - 0.3, 0.22, "主なプロンプト",
                 size=8, bold=True, color=PRIMARY_DARK, font=HEAD_FONT, char_spacing=200)
        add_text(s, x + 0.15, cy + 0.22, cw - 0.3, 1.3, d["prompts"],
                 size=8.5, color=MUTED, font=BODY_FONT, line_spacing=1.45)

    # footer
    add_text(s, 0.6, SH_IN - 0.4, 12.5, 0.3,
             "教材URL：https://atsunaricoda-maker.github.io/claude-chat-course/  "
             "／  6職種：営業 / マーケ・企画 / 人事 / 経理 / 総務・事務 / 管理職",
             size=9, color=MUTED, font=BODY_FONT)


slide_master_timetable()
slide_day_details()

out = "Claude_Chat_5日間研修_カリキュラム表.pptx"
prs.save(out)
print(f"saved: {out}  ({len(prs.slides)} slides)")
